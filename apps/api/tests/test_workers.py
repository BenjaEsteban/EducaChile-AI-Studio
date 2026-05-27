import uuid
from contextlib import contextmanager
from io import BytesIO
from unittest.mock import MagicMock, patch

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from sqlalchemy.orm import configure_mappers

from app.config import settings as app_settings
from app.modules.generation.models import GenerationJob, VideoGenerationSettings
from app.modules.generation.pipeline import (
    _build_narration_chunks,
    _can_reuse_slide_audio_asset,
    _avatar_overlay_type,
    GenerationContext,
    _detect_green_background,
    _ensure_avatar_base_video_asset,
    _load_avatar_source_url,
    _normalize_tts_text,
    resolve_tts_credentials,
    _slide_audio_chunk_specs,
    _slide_segment_duration_seconds,
    compose_segment_for_slide,
    validate_generation_job,
    generate_audio_for_slide,
    generate_avatar_clip_for_slide,
)
from app.modules.jobs.models import Job, JobStatus, JobType
from app.modules.organizations.models import Organization
from app.modules.projects.models import Asset, Presentation, PresentationStatus, Project, ProjectGenerationConfig, Slide
from app.modules.projects.service import MOCK_ORG_ID, MOCK_USER_ID
from app.modules.users.models import User
from app.modules.video.adapters import AvatarVideoProviderError
from app.workers.celery_app import celery_app
from app.workers.tasks import ParsePresentationTask, enqueue_parse_presentation, ping
from tests.conftest import _TestingSession
from tests.fakes import InMemoryStorageProvider
from app.utils.crypto import encrypt_secret

# ── Helpers ───────────────────────────────────────────────────────────────────

def _make_job(job_id: uuid.UUID | None = None) -> Job:
    return Job(
        id=job_id or uuid.uuid4(),
        organization_id=uuid.uuid4(),
        project_id=uuid.uuid4(),
        presentation_id=uuid.uuid4(),
        job_type=JobType.parse_presentation,
        status=JobStatus.queued,
        progress=0.0,
        current_step=None,
        celery_task_id=None,
        error_message=None,
        result=None,
        started_at=None,
        finished_at=None,
    )


def _make_repo(job: Job) -> MagicMock:
    repo = MagicMock()
    repo.get_by_id.return_value = job
    repo.mark_running.return_value = job
    repo.mark_completed.return_value = job
    repo.mark_failed.return_value = job
    repo.update_progress.return_value = job
    return repo


@contextmanager
def _fake_session():
    yield MagicMock()


@contextmanager
def _testing_worker_session():
    db = _TestingSession()
    try:
        yield db
        db.commit()
    except Exception:
        db.rollback()
        raise
    finally:
        db.close()


def _make_pptx_bytes() -> bytes:
    deck = PptxPresentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Clase 1"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255
    )
    slide.placeholders[1].text = "Objetivo de aprendizaje\nContenido visible"

    second = deck.slides.add_slide(deck.slide_layouts[1])
    second.shapes.title.text = "Clase 2"
    second.placeholders[1].text = "Actividad final"

    buffer = BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _create_parse_fixture(storage: InMemoryStorageProvider, pptx_bytes: bytes | None = None):
    storage_key = f"{MOCK_ORG_ID}/{uuid.uuid4()}/deck.pptx"
    storage.upload_file(
        storage_key,
        pptx_bytes if pptx_bytes is not None else _make_pptx_bytes(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    db = _TestingSession()
    try:
        user = User(
            id=MOCK_USER_ID,
            email=f"{uuid.uuid4()}@test.local",
            hashed_password="test",
            full_name="Test User",
        )
        org = Organization(id=MOCK_ORG_ID, name="Test Org", slug=f"test-{uuid.uuid4()}")
        project = Project(
            organization_id=MOCK_ORG_ID,
            owner_id=MOCK_USER_ID,
            name="Project",
        )
        db.add_all([user, org, project])
        db.flush()

        presentation = Presentation(
            project_id=project.id,
            organization_id=MOCK_ORG_ID,
            title="deck.pptx",
            original_filename="deck.pptx",
            storage_key=storage_key,
            status=PresentationStatus.uploaded,
        )
        db.add(presentation)
        db.flush()

        job = Job(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            presentation_id=presentation.id,
            job_type=JobType.parse_presentation,
            status=JobStatus.queued,
        )
        db.add(job)
        db.commit()
        return job.id, presentation.id
    finally:
        db.close()


def test_resolve_tts_credentials_prefers_project_config_and_decrypts(db_session):
    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="TTS Resolution Project",
    )
    db_session.add(project)
    db_session.flush()
    config = ProjectGenerationConfig(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        tts_provider="elevenlabs",
        voice_id="voice-from-project",
        elevenlabs_api_key_encrypted=encrypt_secret("project-secret-key"),
    )
    db_session.add(config)
    db_session.commit()

    resolution = resolve_tts_credentials(config, app_settings)

    assert resolution.provider == "elevenlabs"
    assert resolution.api_key == "project-secret-key"
    assert resolution.voice_id == "voice-from-project"
    assert resolution.credentials_source == "project_config"


def test_resolve_tts_credentials_uses_env_fallback(monkeypatch):
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.TTS_PROVIDER", "elevenlabs")
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.ELEVENLABS_API_KEY", "env-secret-key")
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.ELEVENLABS_VOICE_ID", "env-voice-id")

    resolution = resolve_tts_credentials(None, app_settings)

    assert resolution.provider == "elevenlabs"
    assert resolution.api_key == "env-secret-key"
    assert resolution.voice_id == "env-voice-id"
    assert resolution.credentials_source == "env_fallback"


def test_validate_generation_job_logs_safe_tts_resolution(caplog, db_session, monkeypatch):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Safe Log Project",
    )
    db_session.add(project)
    db_session.flush()
    config = ProjectGenerationConfig(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        tts_provider="elevenlabs",
        voice_id="voice-from-project",
        elevenlabs_api_key_encrypted=encrypt_secret("project-secret-key"),
    )
    db_session.add(config)
    db_session.flush()
    video_settings = VideoGenerationSettings(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        validation_status="saved",
        wavespeed_api_key_encrypted=encrypt_secret("wavespeed-secret"),
        avatar_source_url="https://public.example.test/avatar.png",
    )
    db_session.add(video_settings)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    db_session.add(
        Slide(
            presentation_id=presentation.id,
            position=1,
            title="Slide 1",
            notes="Narration text",
            thumbnail_key=avatar_key,
            metadata_={"dialogue": "Narration text", "rendered_image_key": avatar_key},
        )
    )
    db_session.commit()
    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.commit()

    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.WAVESPEED_API_KEY", "wavespeed-env-key")
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.ALLOW_DUMMY_TTS", False)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)
    monkeypatch.setattr("app.modules.generation.pipeline._load_slide_preview_image", lambda *_args, **_kwargs: (b"slide-preview", {"storage_key": avatar_key, "includes_text": True, "asset_type": "slide_preview"}))

    with caplog.at_level("INFO"):
        context = validate_generation_job(db_session, job, project.id, MOCK_ORG_ID, storage)

    assert context.elevenlabs_api_key == "project-secret-key"
    assert context.elevenlabs_voice_id == "voice-from-project"
    assert "project-secret-key" not in caplog.text
    assert "voice-from-project" not in caplog.text
    assert "credentials_source=project_config" in caplog.text


# ── ping ──────────────────────────────────────────────────────────────────────

def test_worker_imports_register_all_sqlalchemy_mappers():
    configure_mappers()


def test_ping_returns_message():
    celery_app.conf.task_always_eager = True
    result = ping.apply(kwargs={"message": "hello"})
    assert result.result["message"] == "hello"


def test_ping_default_message():
    celery_app.conf.task_always_eager = True
    result = ping.apply()
    assert result.result["message"] == "pong"


# ── JobTask base — ciclo de vida ──────────────────────────────────────────────

def test_job_task_marks_running_then_completed():
    from app.workers.base_task import JobTask

    class EchoTask(JobTask):
        name = "test.echo"
        def run_job(self, job_id, payload="ok", **kwargs):
            return {"echo": payload}

    job = _make_job()
    repo = _make_repo(job)

    with patch("app.workers.base_task.worker_db_session", _fake_session), \
         patch("app.workers.base_task.JobRepository", return_value=repo):
        task = celery_app.register_task(EchoTask())
        celery_app.conf.task_always_eager = True
        result = task.apply(kwargs={"job_id": str(job.id), "payload": "test"})

    assert result.result == {"echo": "test"}
    repo.mark_running.assert_called_once()
    repo.mark_completed.assert_called_once()
    repo.mark_failed.assert_not_called()


def test_job_task_marks_failed_on_exception():
    from app.workers.base_task import JobTask

    class FailTask(JobTask):
        name = "test.fail"
        def run_job(self, job_id, **kwargs):
            raise ValueError("error simulado")

    job = _make_job()
    repo = _make_repo(job)

    with patch("app.workers.base_task.worker_db_session", _fake_session), \
         patch("app.workers.base_task.JobRepository", return_value=repo):
        task = celery_app.register_task(FailTask())
        celery_app.conf.task_always_eager = True
        result = task.apply(kwargs={"job_id": str(job.id)})

    # En modo eager, las excepciones se capturan en el result
    assert result.failed()
    repo.mark_running.assert_called_once()
    repo.mark_failed.assert_called_once()
    _, kwargs = repo.mark_failed.call_args
    assert "error simulado" in kwargs["error_message"]


def test_job_task_job_not_found_returns_error():
    from app.workers.base_task import JobTask

    class AnyTask(JobTask):
        name = "test.any"
        def run_job(self, job_id, **kwargs):
            return {}

    repo = MagicMock()
    repo.get_by_id.return_value = None

    with patch("app.workers.base_task.worker_db_session", _fake_session), \
         patch("app.workers.base_task.JobRepository", return_value=repo):
        task = celery_app.register_task(AnyTask())
        celery_app.conf.task_always_eager = True
        result = task.apply(kwargs={"job_id": str(uuid.uuid4())})

    assert result.result == {"error": "job_not_found"}
    repo.mark_running.assert_not_called()


def test_job_task_set_progress_updates_db():
    from app.workers.base_task import JobTask

    class ProgressTask(JobTask):
        name = "test.progress"
        def run_job(self, job_id, **kwargs):
            self.set_progress(job_id, 50.0)
            self.set_progress(job_id, 100.0)
            return {"done": True}

    job = _make_job()
    repo = _make_repo(job)

    with patch("app.workers.base_task.worker_db_session", _fake_session), \
         patch("app.workers.base_task.JobRepository", return_value=repo):
        task = celery_app.register_task(ProgressTask())
        celery_app.conf.task_always_eager = True
        task.apply(kwargs={"job_id": str(job.id)})

    assert repo.update_progress.call_count == 2
    calls = repo.update_progress.call_args_list
    assert calls[0].args[1] == 50.0
    assert calls[1].args[1] == 100.0


# ── ParsePresentationTask ─────────────────────────────────────────────────────

def test_parse_presentation_completes_successfully():
    storage = InMemoryStorageProvider()
    job_id, presentation_id = _create_parse_fixture(storage)

    with patch("app.workers.base_task.worker_db_session", _testing_worker_session), \
         patch("app.workers.tasks.worker_db_session", _testing_worker_session), \
         patch("app.workers.tasks.get_storage", return_value=storage):
        celery_app.conf.task_always_eager = True
        result = ParsePresentationTask().apply(kwargs={
            "job_id": str(job_id),
            "presentation_id": str(presentation_id),
        })

    assert not result.failed()
    body = result.result
    assert body["parsed"] is True
    assert body["presentation_id"] == str(presentation_id)
    assert body["slide_count"] == 2

    db = _TestingSession()
    try:
        presentation = db.get(Presentation, presentation_id)
        slides = (
            db.query(Slide)
            .filter(Slide.presentation_id == presentation_id)
            .order_by(Slide.position)
            .all()
        )
        job = db.get(Job, job_id)
        assert presentation.status == PresentationStatus.parsed
        assert presentation.slide_count == 2
        assert job.status == JobStatus.completed
        assert job.progress == 100.0
        assert job.current_step == "Completed"
        assert len(slides) == 2
        assert slides[0].position == 1
        assert slides[0].title == "Clase 1"
        assert "Objetivo de aprendizaje" in slides[0].metadata_["visible_text"]
        assert slides[0].notes is None
        assert slides[0].metadata_["dialogue"] == ""
    finally:
        db.close()


def test_parse_presentation_saves_thumbnails_when_preview_rendering_succeeds():
    storage = InMemoryStorageProvider()
    job_id, presentation_id = _create_parse_fixture(storage)

    with patch("app.workers.base_task.worker_db_session", _testing_worker_session), \
         patch("app.workers.tasks.worker_db_session", _testing_worker_session), \
         patch("app.workers.tasks.get_storage", return_value=storage), \
         patch(
             "app.workers.tasks.render_slide_previews",
             return_value={1: "presentations/test/previews/slide-1.png", 2: "presentations/test/previews/slide-2.png"},
         ):
        celery_app.conf.task_always_eager = True
        result = ParsePresentationTask().apply(kwargs={
            "job_id": str(job_id),
            "presentation_id": str(presentation_id),
        })

    assert not result.failed()

    db = _TestingSession()
    try:
        slides = (
            db.query(Slide)
            .filter(Slide.presentation_id == presentation_id)
            .order_by(Slide.position)
            .all()
        )
        assert slides[0].thumbnail_key == "presentations/test/previews/slide-1.png"
        assert slides[0].metadata_["rendered_image_key"] == "presentations/test/previews/slide-1.png"
        assert slides[0].metadata_["slide_preview"]["asset_type"] == "slide_preview"
        assert slides[0].metadata_["slide_preview"]["includes_text"] is True
    finally:
        db.close()


def test_parse_presentation_failure_marks_presentation_failed():
    storage = InMemoryStorageProvider()
    job_id, presentation_id = _create_parse_fixture(storage, pptx_bytes=b"not a pptx")

    with patch("app.workers.base_task.worker_db_session", _testing_worker_session), \
         patch("app.workers.tasks.worker_db_session", _testing_worker_session), \
         patch("app.workers.tasks.get_storage", return_value=storage):
        celery_app.conf.task_always_eager = True
        result = ParsePresentationTask().apply(kwargs={
            "job_id": str(job_id),
            "presentation_id": str(presentation_id),
        })

    assert result.failed()
    db = _TestingSession()
    try:
        presentation = db.get(Presentation, presentation_id)
        job = db.get(Job, job_id)
        assert presentation.status == PresentationStatus.failed
        assert job.status == JobStatus.failed
        assert job.error_message
        assert job.current_step == "Failed"
    finally:
        db.close()


# ── enqueue_parse_presentation ────────────────────────────────────────────────

def test_enqueue_returns_celery_task_id():
    mock_result = MagicMock()
    mock_result.id = "fake-celery-id"

    with patch("app.workers.tasks.parse_presentation.apply_async", return_value=mock_result):
        task_id = enqueue_parse_presentation(uuid.uuid4(), uuid.uuid4())

    assert task_id == "fake-celery-id"


def test_enqueue_uses_presentations_queue():
    mock_result = MagicMock()
    mock_result.id = "x"

    with patch(
        "app.workers.tasks.parse_presentation.apply_async",
        return_value=mock_result,
    ) as mock_apply:
        enqueue_parse_presentation(uuid.uuid4(), uuid.uuid4())

    assert mock_apply.call_args[1]["queue"] == "presentations"


def test_enqueue_passes_job_and_presentation_ids():
    job_id = uuid.uuid4()
    presentation_id = uuid.uuid4()
    mock_result = MagicMock()
    mock_result.id = "x"

    with patch(
        "app.workers.tasks.parse_presentation.apply_async",
        return_value=mock_result,
    ) as mock_apply:
        enqueue_parse_presentation(job_id, presentation_id)

    kwargs = mock_apply.call_args[1]["kwargs"]
    assert kwargs["job_id"] == str(job_id)
    assert kwargs["presentation_id"] == str(presentation_id)


def test_generate_avatar_clip_uses_audio_lipsync_flow(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")
    slide_preview_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_preview_key, b"slide-preview", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Talking Photo Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_preview_key,
        metadata_={
            "dialogue": "Narration text for talking photo",
            "rendered_image_key": slide_preview_key,
            "slide_preview": {
                "asset_type": "slide_preview",
                "storage_key": slide_preview_key,
                "render_source": "ppt_render",
                "includes_text": True,
            },
        },
    )
    db_session.add(slide)
    db_session.flush()
    db_session.add(
        Asset(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            slide_id=None,
            asset_type="avatar_source",
            storage_key=avatar_key,
            filename="avatar.png",
            mime_type="image/png",
            size_bytes=len(b"avatar-bytes"),
        )
    )
    base_video_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="avatar_base_video",
        storage_key=f"projects/{project.id}/avatar/base/avatar-base.mp4",
        filename="avatar-base.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"base-video-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "generated_from": "real_avatar_base_video",
            "fallback_used": False,
            "base_video_provider": "wavespeed_infinitetalk_fast",
            "avatar_source_signature": "signature",
        },
    )
    db_session.add(base_video_asset)
    db_session.flush()
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=5.0,
        metadata_json={"generation_job_id": "job-123", "slide_position": 1},
    )
    db_session.add(audio_asset)
    db_session.commit()

    class FakeResponse:
        def __init__(self, content=b"video-bytes", status_code=200):
            self.content = content
            self.status_code = status_code

    class FakeProvider:
        def generate_avatar_video_from_audio(
            self,
            image_url,
            audio_url,
            duration=5,
            seed=-1,
            prompt=None,
            resolution=None,
            api_key=None,
            audio_duration_seconds=None,
            image_bytes=None,
            audio_bytes=None,
            image_filename=None,
            image_content_type=None,
            audio_filename=None,
            audio_content_type=None,
            retry_on_mismatch=True,
            minimum_duration_ratio=0.8,
            heartbeat_callback=None,
        ):
            assert image_bytes == b"avatar-bytes"
            assert audio_bytes == b"audio-bytes"
            assert image_filename == "avatar.png"
            assert image_content_type == "image/png"
            assert audio_filename == "slide-1-chunk-1.mp3"
            assert audio_content_type == "audio/mpeg"
            assert prompt is None
            assert duration == 5
            assert resolution == "480p"
            assert audio_duration_seconds == 5.0
            assert retry_on_mismatch is True
            assert minimum_duration_ratio == 0.8
            assert heartbeat_callback is not None
            self.last_image_url = "https://wavespeed.test/uploaded-image.png"
            self.last_audio_url = "https://wavespeed.test/uploaded-audio.mp3"
            self.last_external_checks = {
                "image": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/uploaded-image.png"},
                "audio": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/uploaded-audio.mp3"},
            }
            self.last_duration_ratio = 1.0
            self.last_request_id = "request-123"
            return "https://cdn.test/out.mp4"

    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_GENERATION_MODE", "infinitetalk_image")
    monkeypatch.setattr(
        "app.modules.generation.pipeline.settings.AVATAR_LIPSYNC_PROVIDER",
        "wavespeed_infinitetalk",
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline.get_avatar_video_provider",
        lambda *_args, **_kwargs: FakeProvider(),
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline._talking_photo_duration_from_audio",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(
            AssertionError("audio_lipsync must not use talking-photo duration logic")
        ),
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline.httpx.get",
        lambda url, timeout=None: FakeResponse(),
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline.httpx.head",
        lambda url, timeout=None, follow_redirects=True: FakeResponse(status_code=200),
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline.ComposerService.normalize_audio_to_mp3",
        lambda self, audio_bytes: b"normalized-audio",
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline.ComposerService.strip_audio_from_video",
        lambda self, video_bytes: b"avatar-video-bytes",
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline._probe_media_info",
        lambda media_bytes, *_args, **_kwargs: (
            {
                "duration_seconds": 5.0,
                "has_video": True,
                "has_audio": True,
                "video_codec": "h264",
                "audio_codec": "aac",
                "width": 1280,
                "height": 720,
            }
            if media_bytes == b"video-bytes"
            else {
                "duration_seconds": 5.0,
                "has_video": False,
                "has_audio": True,
                "video_codec": None,
                "audio_codec": "aac",
                "width": None,
                "height": None,
            }
            if media_bytes in {b"audio-bytes", b"normalized-audio"}
            else {
                "duration_seconds": 5.0,
                "has_video": True,
                "has_audio": False,
                "video_codec": "h264",
                "audio_codec": None,
                "width": 1280,
                "height": 720,
            }
            if media_bytes == b"avatar-video-bytes"
            else {
                "duration_seconds": 5.0,
                "has_video": False,
                "has_audio": False,
                "video_codec": None,
                "audio_codec": None,
                "width": None,
                "height": None,
            }
        ),
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline._analyze_video_motion",
        lambda *_args, **_kwargs: {
            "motion_score": 5.0,
            "almost_static": False,
            "sample_count": 3,
        },
    )
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_avatar_clip_for_slide(
        db_session,
        storage,
        job,
        context,
        slide,
        1,
        1,
        audio_asset,
    )

    assert asset.asset_type == "generated_avatar_clip"
    assert asset.metadata_json["mode"] == "infinitetalk_image"
    assert asset.metadata_json["selected_provider"] == "wavespeed_infinitetalk"
    assert asset.metadata_json["provider"] == "wavespeed"
    assert asset.metadata_json["wavespeed_request_id"] == "request-123"
    assert asset.metadata_json["source_audio_url"] == "https://wavespeed.test/uploaded-audio.mp3"
    assert asset.metadata_json["chunks"][0]["text"] == "Narration text for talking photo."
    assert asset.metadata_json["chunks"][0]["audio_url"] == "https://wavespeed.test/uploaded-audio.mp3"
    assert asset.metadata_json["chunks"][0]["measured_tts_duration"] == 5.0
    assert asset.metadata_json["ffprobe"]["has_video"] is True
    assert asset.metadata_json["ffprobe"]["has_audio"] is False
    assert asset.metadata_json["chunk_count"] == 1
    assert asset.metadata_json["chunks"][0]["provider_audio_present"] is True
    assert asset.metadata_json["chunks"][0]["stripped_clip_ffprobe"]["has_audio"] is False
    assert asset.metadata_json["chunks"][0]["image_url"] == "https://wavespeed.test/uploaded-image.png"
    assert asset.metadata_json["chunks"][0]["image_url_external_check_result"]["validated"] is True
    assert asset.metadata_json["chunks"][0]["audio_url_external_check_result"]["validated"] is True


def test_generate_avatar_clip_uses_fast_lipsync_base_video(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")
    slide_preview_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_preview_key, b"slide-preview", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Fast Lipsync Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_preview_key,
        metadata_={
            "dialogue": "Narration text for fast lipsync",
            "rendered_image_key": slide_preview_key,
            "slide_preview": {
                "asset_type": "slide_preview",
                "storage_key": slide_preview_key,
                "render_source": "ppt_render",
                "includes_text": True,
            },
        },
    )
    db_session.add(slide)
    db_session.flush()
    db_session.add(
        Asset(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            slide_id=None,
            asset_type="avatar_source",
            storage_key=avatar_key,
            filename="avatar.png",
            mime_type="image/png",
            size_bytes=len(b"avatar-bytes"),
        )
    )
    base_video_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="avatar_base_video",
        storage_key=f"projects/{project.id}/avatar/base/avatar-base.mp4",
        filename="avatar-base.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"base-video-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "avatar_base_video_source": "generated_real_avatar_base_video",
            "avatar_base_video_is_real_motion": True,
            "fallback_used": False,
            "base_video_provider": "wavespeed_infinitetalk_fast",
        },
    )
    db_session.add(base_video_asset)
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "generation_job_id": "job-123",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "text": "Narration text for fast lipsync",
                    "word_count": 4,
                    "estimated_duration_seconds": 6.0,
                    "measured_tts_duration": 6.0,
                    "audio_storage_key": audio_key,
                    "audio_url": "https://wavespeed.test/uploaded-audio.mp3",
                }
            ],
        },
    )
    db_session.add(audio_asset)
    db_session.commit()

    class FakeResponse:
        def __init__(self, content=b"video-bytes", status_code=200):
            self.content = content
            self.status_code = status_code

    class FakeProvider:
        def generate_avatar_video_from_base_video(
            self,
            *,
            base_video_url,
            audio_url,
            duration=5,
            seed=-1,
            api_key=None,
            audio_duration_seconds=None,
            base_video_bytes=None,
            base_video_filename=None,
            base_video_content_type=None,
            audio_bytes=None,
            audio_filename=None,
            audio_content_type=None,
            sync_mode=None,
            model_name=None,
            retry_on_mismatch=True,
            minimum_duration_ratio=0.8,
            heartbeat_callback=None,
        ):
            assert base_video_bytes == b"base-video-bytes"
            assert audio_bytes == b"audio-bytes"
            assert base_video_filename == "avatar-base.mp4"
            assert audio_filename == "slide-1-chunk-1.mp3"
            assert base_video_content_type == "video/mp4"
            assert audio_content_type == "audio/mpeg"
            assert sync_mode == "loop"
            assert model_name == "wavespeed-ai/sync-lipsync-3"
            assert duration == 6
            assert audio_duration_seconds == 6.0
            assert retry_on_mismatch is True
            assert minimum_duration_ratio == 0.8
            assert heartbeat_callback is not None
            self.last_image_url = "https://wavespeed.test/base-uploaded.mp4"
            self.last_audio_url = "https://wavespeed.test/audio-uploaded.mp3"
            self.last_external_checks = {
                "video": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/base-uploaded.mp4"},
                "audio": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/audio-uploaded.mp3"},
            }
            self.last_duration_ratio = 1.0
            self.last_request_id = "request-fast-123"
            return "https://cdn.test/fast.mp4"

    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_GENERATION_MODE", "fast_lipsync")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_LIPSYNC_PROVIDER", "wavespeed_sync_lipsync_3")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_SYNC_MODE", "loop")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_LIPSYNC_MODEL_PATH", "wavespeed-ai/sync-lipsync-3")
    monkeypatch.setattr("app.modules.generation.pipeline.get_avatar_video_provider", lambda *_args, **_kwargs: FakeProvider())
    monkeypatch.setattr(
        "app.modules.generation.pipeline._ensure_avatar_base_video_asset",
        lambda **_kwargs: (base_video_asset, b"base-video-bytes", base_video_asset.metadata_json or {}),
    )
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: (
        {
            "duration_seconds": 6.0,
            "has_video": True,
            "has_audio": False,
            "video_codec": "h264",
            "audio_codec": None,
            "width": 1280,
            "height": 720,
        }
        if media_bytes == b"base-video-bytes"
        else {
            "duration_seconds": 6.0,
            "has_video": False,
            "has_audio": True,
            "video_codec": None,
            "audio_codec": "aac",
            "width": None,
            "height": None,
        }
        if media_bytes == b"audio-bytes"
        else {
            "duration_seconds": 6.0,
            "has_video": True,
            "has_audio": False,
            "video_codec": "h264",
            "audio_codec": None,
            "width": 1280,
            "height": 720,
        }
    ))
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.get", lambda url, timeout=None: FakeResponse())
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.head", lambda url, timeout=None, follow_redirects=True: FakeResponse(status_code=200))
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.strip_audio_from_video", lambda self, video_bytes: video_bytes)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_avatar_clip_for_slide(
        db_session,
        storage,
        job,
        context,
        slide,
        1,
        1,
        audio_asset,
    )

    assert asset.asset_type == "generated_avatar_clip"
    assert asset.metadata_json["mode"] == "fast_lipsync"
    assert asset.metadata_json["selected_provider"] == "wavespeed_sync_lipsync_3"
    assert asset.metadata_json["provider"] == "wavespeed"
    assert asset.metadata_json["wavespeed_request_id"] == "request-fast-123"
    assert asset.metadata_json["source_audio_url"] == "https://wavespeed.test/audio-uploaded.mp3"
    assert asset.metadata_json["source_image_url"] == "https://wavespeed.test/base-uploaded.mp4"
    assert asset.metadata_json["avatar_base_video_source"] == "generated_real_avatar_base_video"
    assert asset.metadata_json["avatar_base_video_is_real_motion"] is True
    assert asset.metadata_json["avatar_overlay_type"] == "provider_lipsync_video"
    assert asset.metadata_json["provider_lipsync_output_present"] is True
    assert asset.metadata_json["final_overlay_source"] == "provider_lipsync_output"
    assert asset.metadata_json["chunks"][0]["fallback_used"] is False


def test_generate_avatar_clip_uses_image_audio_infinitetalk_mode(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")
    slide_preview_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_preview_key, b"slide-preview", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Image Audio Mode Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_preview_key,
        metadata_={
            "dialogue": "Narration text for image audio mode",
            "rendered_image_key": slide_preview_key,
            "slide_preview": {
                "asset_type": "slide_preview",
                "storage_key": slide_preview_key,
                "render_source": "ppt_render",
                "includes_text": True,
            },
        },
    )
    db_session.add(slide)
    db_session.flush()
    db_session.add(
        Asset(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            slide_id=None,
            asset_type="avatar_source",
            storage_key=avatar_key,
            filename="avatar.png",
            mime_type="image/png",
            size_bytes=len(b"avatar-bytes"),
        )
    )
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "generation_job_id": "job-123",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "text": "Narration text for image audio mode",
                    "word_count": 5,
                    "estimated_duration_seconds": 6.0,
                    "measured_tts_duration": 6.0,
                    "audio_storage_key": audio_key,
                    "audio_url": "https://wavespeed.test/uploaded-audio.mp3",
                }
            ],
        },
    )
    db_session.add(audio_asset)
    db_session.commit()

    class FakeResponse:
        def __init__(self, content=b"video-bytes", status_code=200):
            self.content = content
            self.status_code = status_code

    class FakeProvider:
        def generate_avatar_video_from_audio(self, **kwargs):
            assert kwargs["image_bytes"] == b"avatar-bytes"
            assert kwargs["audio_bytes"] == b"audio-bytes"
            assert kwargs["audio_url"] == "https://wavespeed.test/uploaded-audio.mp3"
            assert kwargs["resolution"] == "480p"
            self.last_request_id = "request-image-audio-123"
            self.last_image_url = "https://wavespeed.test/uploaded-image.mp4"
            self.last_audio_url = "https://wavespeed.test/uploaded-audio.mp3"
            self.last_external_checks = {
                "image": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/uploaded-image.mp4"},
                "audio": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/uploaded-audio.mp3"},
            }
            self.last_status_history = [{"status": "completed", "elapsed_seconds": 2.0, "outputs_present": True}]
            self.last_duration_ratio = 1.0
            return "https://cdn.test/image-audio.mp4"

        def generate_avatar_video_from_base_video(self, **_kwargs):
            raise AssertionError("image_audio_infinitetalk must not use base video")

    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_GENERATION_MODE", "image_audio_infinitetalk")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_IMAGE_AUDIO_PROVIDER", "wavespeed_infinitetalk_fast")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_IMAGE_AUDIO_RESOLUTION", "480p")
    monkeypatch.setattr("app.modules.generation.pipeline.get_avatar_video_provider", lambda *_args, **_kwargs: FakeProvider())
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.get", lambda url, timeout=None: FakeResponse())
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.head", lambda url, timeout=None, follow_redirects=True: FakeResponse(status_code=200))
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    } if media_bytes == b"video-bytes" else {
        "duration_seconds": 6.0,
        "has_video": False,
        "has_audio": True,
        "video_stream_count": 0,
        "audio_stream_count": 1,
        "video_codec": None,
        "audio_codec": "aac",
        "width": None,
        "height": None,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.strip_audio_from_video", lambda self, video_bytes: video_bytes)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_avatar_clip_for_slide(
        db_session,
        storage,
        job,
        context,
        slide,
        1,
        1,
        audio_asset,
    )

    assert asset.metadata_json["mode"] == "image_audio_infinitetalk"
    assert asset.metadata_json["avatar_overlay_type"] == "image_audio_infinitetalk_video"
    assert asset.metadata_json["provider_request_type"] == "image_plus_audio"
    assert asset.metadata_json["avatar_provider_name"] == "wavespeed_infinitetalk_fast"
    assert asset.metadata_json["provider_lipsync_output_present"] is True
    assert asset.metadata_json["final_overlay_source"] == "provider_lipsync_output"
    assert asset.metadata_json["avatar_base_video_source"] is None


def test_compose_segment_uses_provider_overlay_without_base_video_metadata(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    slide_image_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_image_key, b"slide-image", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.mp4"
    storage.upload_file(avatar_key, b"avatar-video-bytes", "video/mp4")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Compose Provider Overlay Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_image_key,
        metadata_={
            "rendered_image_key": slide_image_key,
            "slide_preview": {
                "asset_type": "slide_preview",
                "storage_key": slide_image_key,
                "render_source": "ppt_render",
                "includes_text": True,
            },
        },
    )
    db_session.add(slide)
    db_session.flush()
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "generation_job_id": "job-123",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "text": "Narration text",
                    "word_count": 2,
                    "estimated_duration_seconds": 6.0,
                    "measured_tts_duration": 6.0,
                    "audio_storage_key": audio_key,
                    "audio_url": "https://wavespeed.test/uploaded-audio.mp3",
                }
            ],
        },
    )
    avatar_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="generated_avatar_clip",
        storage_key=avatar_key,
        filename="avatar-slide-1.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"avatar-video-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "avatar_overlay_type": "provider_lipsync_video",
            "fallback_used": False,
            "fallback_reason": None,
            "provider_name": "wavespeed",
            "provider_prediction_id": "pred-1",
            "provider_audio_present": False,
            "motion_analysis": {"almost_static": False},
            "avatar_video_storage_key": avatar_key,
            "avatar_video_duration": 6.0,
            "avatar_video_has_motion_checked": True,
            "source_audio_url": "https://wavespeed.test/uploaded-audio.mp3",
        },
    )
    db_session.add_all([audio_asset, avatar_asset])
    db_session.commit()

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    compose_calls = {}

    def fake_compose(self, **kwargs):
        compose_calls.update(kwargs)
        return b"segment-bytes"

    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.compose_slide_video", fake_compose)
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": True,
        "video_stream_count": 1,
        "audio_stream_count": 1,
        "video_codec": "h264",
        "audio_codec": "aac",
        "width": 1280,
        "height": 720,
    } if media_bytes == b"segment-bytes" else {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    asset = compose_segment_for_slide(db_session, storage, job, context, slide, 1, 1, audio_asset, avatar_asset)

    assert asset.metadata_json["composition_used_generated_avatar_clip"] is True
    assert asset.metadata_json["avatar_overlay_type"] == "provider_lipsync_video"
    assert compose_calls["avatar_chromakey"] is False


def test_compose_segment_static_fallback_only_fails_when_strict(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    slide_image_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_image_key, b"slide-image", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.mp4"
    storage.upload_file(avatar_key, b"avatar-video-bytes", "video/mp4")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Compose Static Fallback Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_image_key,
        metadata_={"rendered_image_key": slide_image_key, "slide_preview": {"asset_type": "slide_preview", "storage_key": slide_image_key, "render_source": "ppt_render", "includes_text": True}},
    )
    db_session.add(slide)
    db_session.flush()
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
    )
    avatar_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="generated_avatar_clip",
        storage_key=avatar_key,
        filename="avatar-slide-1.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"avatar-video-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "avatar_overlay_type": "static_avatar_fallback",
            "fallback_used": True,
            "fallback_reason": "provider timeout",
            "provider_name": "wavespeed",
            "provider_prediction_id": None,
            "avatar_video_storage_key": avatar_key,
            "avatar_video_duration": 6.0,
            "avatar_video_has_motion_checked": True,
        },
    )
    db_session.add_all([audio_asset, avatar_asset])
    db_session.commit()

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    monkeypatch.setattr("app.modules.generation.pipeline.settings.FAIL_ON_STATIC_AVATAR_FALLBACK", True)
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.compose_slide_video", lambda self, **kwargs: b"segment-bytes")
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": True,
        "video_stream_count": 1,
        "audio_stream_count": 1,
        "video_codec": "h264",
        "audio_codec": "aac",
        "width": 1280,
        "height": 720,
    } if media_bytes == b"segment-bytes" else {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    try:
        compose_segment_for_slide(db_session, storage, job, context, slide, 1, 1, audio_asset, avatar_asset)
    except Exception as exc:
        assert "static fallback" in str(exc)
    else:
        raise AssertionError("strict static fallback should fail")


def test_compose_segment_provider_green_output_uses_cleanup_when_global_chromakey_is_disabled(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    slide_image_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_image_key, b"slide-image", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.mp4"
    storage.upload_file(avatar_key, b"avatar-video-bytes", "video/mp4")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Chromakey Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_image_key,
        metadata_={"rendered_image_key": slide_image_key, "slide_preview": {"asset_type": "slide_preview", "storage_key": slide_image_key, "render_source": "ppt_render", "includes_text": True}},
    )
    db_session.add(slide)
    db_session.flush()
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
    )
    avatar_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="generated_avatar_clip",
        storage_key=avatar_key,
        filename="avatar-slide-1.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"avatar-video-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "avatar_overlay_type": "provider_lipsync_video",
            "fallback_used": False,
            "provider_name": "wavespeed",
            "provider_prediction_id": "pred-1",
            "avatar_video_storage_key": avatar_key,
            "avatar_video_duration": 6.0,
            "avatar_video_has_motion_checked": True,
            "green_screen_background": True,
        },
    )
    db_session.add_all([audio_asset, avatar_asset])
    db_session.commit()

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    compose_calls = {}

    def fake_compose(self, **kwargs):
        compose_calls.update(kwargs)
        return b"segment-bytes"

    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.compose_slide_video", fake_compose)
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": True,
        "video_stream_count": 1,
        "audio_stream_count": 1,
        "video_codec": "h264",
        "audio_codec": "aac",
        "width": 1280,
        "height": 720,
    } if media_bytes == b"segment-bytes" else {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)
    monkeypatch.setattr(
        "app.modules.generation.pipeline._detect_green_background",
        lambda _bytes: {"detected": True, "green_ratio": 0.91, "sample_count": 1024, "threshold": 0.45},
    )

    asset = compose_segment_for_slide(db_session, storage, job, context, slide, 1, 1, audio_asset, avatar_asset)
    assert asset.duration_seconds == 6.0
    assert compose_calls["avatar_chromakey"] is True
    assert asset.metadata_json["provider_requires_chromakey"] is True
    assert asset.metadata_json["green_background_cleanup_applied"] is True


def test_compose_segment_fast_lipsync_rejects_missing_provider_output(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    slide_image_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_image_key, b"slide-image", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar-base.mp4"
    storage.upload_file(avatar_key, b"avatar-base-video", "video/mp4")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Fast Lipsync Missing Provider Output",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_image_key,
        metadata_={"rendered_image_key": slide_image_key, "slide_preview": {"asset_type": "slide_preview", "storage_key": slide_image_key, "render_source": "ppt_render", "includes_text": True}},
    )
    db_session.add(slide)
    db_session.flush()
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
    )
    avatar_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="generated_avatar_clip",
        storage_key=avatar_key,
        filename="avatar-base.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"avatar-base-video"),
        duration_seconds=4.0,
        metadata_json={
            "mode": "fast_lipsync",
            "avatar_overlay_type": "provider_lipsync_video",
            "provider_lipsync_output_present": False,
            "final_overlay_source": "static_avatar_fallback",
            "fallback_used": False,
            "provider_name": "wavespeed",
            "provider_prediction_id": "pred-1",
            "avatar_video_storage_key": avatar_key,
            "avatar_video_duration": 4.0,
            "avatar_video_has_motion_checked": True,
            "avatar_base_video_source": "generated_real_avatar_base_video",
            "avatar_base_video_is_real_motion": True,
        },
    )
    db_session.add_all([audio_asset, avatar_asset])
    db_session.commit()

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    monkeypatch.setattr("app.modules.generation.pipeline.settings.FAIL_ON_STATIC_AVATAR_FALLBACK", False)
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.compose_slide_video", lambda self, **kwargs: b"segment-bytes")
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": True,
        "video_stream_count": 1,
        "audio_stream_count": 1,
        "video_codec": "h264",
        "audio_codec": "aac",
        "width": 1280,
        "height": 720,
    } if media_bytes == b"segment-bytes" else {
        "duration_seconds": 4.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    try:
        compose_segment_for_slide(db_session, storage, job, context, slide, 1, 1, audio_asset, avatar_asset)
    except Exception as exc:
        assert "per-slide lip-sync output" in str(exc)
    else:
        raise AssertionError("fast_lipsync must not compose without a provider output")


def test_compose_segment_chromakey_applies_only_when_enabled(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    slide_image_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_image_key, b"slide-image", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.mp4"
    storage.upload_file(avatar_key, b"avatar-video-bytes", "video/mp4")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Chromakey Enabled Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_image_key,
        metadata_={"rendered_image_key": slide_image_key, "slide_preview": {"asset_type": "slide_preview", "storage_key": slide_image_key, "render_source": "ppt_render", "includes_text": True}},
    )
    db_session.add(slide)
    db_session.flush()
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
    )
    avatar_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="generated_avatar_clip",
        storage_key=avatar_key,
        filename="avatar-slide-1.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"avatar-video-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "avatar_overlay_type": "provider_lipsync_video",
            "fallback_used": False,
            "provider_name": "wavespeed",
            "provider_prediction_id": "pred-1",
            "avatar_video_storage_key": avatar_key,
            "avatar_video_duration": 6.0,
            "avatar_video_has_motion_checked": True,
            "green_screen_background": True,
        },
    )
    db_session.add_all([audio_asset, avatar_asset])
    db_session.commit()

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    compose_calls = {}

    def fake_compose(self, **kwargs):
        compose_calls.update(kwargs)
        return b"segment-bytes"

    monkeypatch.setattr("app.modules.generation.pipeline.settings.ENABLE_AVATAR_CHROMAKEY", True)
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.compose_slide_video", fake_compose)
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": True,
        "video_stream_count": 1,
        "audio_stream_count": 1,
        "video_codec": "h264",
        "audio_codec": "aac",
        "width": 1280,
        "height": 720,
    } if media_bytes == b"segment-bytes" else {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)
    monkeypatch.setattr(
        "app.modules.generation.pipeline._detect_green_background",
        lambda _bytes: {"detected": True, "green_ratio": 0.91, "sample_count": 1024, "threshold": 0.45},
    )

    asset = compose_segment_for_slide(db_session, storage, job, context, slide, 1, 1, audio_asset, avatar_asset)
    assert asset.duration_seconds == 6.0
    assert compose_calls["avatar_chromakey"] is True
    assert asset.metadata_json["provider_requires_chromakey"] is True
    assert asset.metadata_json["green_background_cleanup_applied"] is True


def test_generate_avatar_clip_falls_back_to_static_avatar_when_provider_times_out(
    monkeypatch,
    db_session,
):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")
    base_video_key = f"projects/{uuid.uuid4()}/avatar/avatar-base.mp4"
    storage.upload_file(base_video_key, b"base-video-bytes", "video/mp4")
    slide_preview_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_preview_key, b"slide-preview", "image/png")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Fast Lipsync Fallback Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_preview_key,
        metadata_={
            "dialogue": "Narration text for fallback",
            "rendered_image_key": slide_preview_key,
            "slide_preview": {
                "asset_type": "slide_preview",
                "storage_key": slide_preview_key,
                "render_source": "ppt_render",
                "includes_text": True,
            },
        },
    )
    db_session.add(slide)
    db_session.flush()
    db_session.add(
        Asset(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            slide_id=None,
            asset_type="avatar_source",
            storage_key=avatar_key,
            filename="avatar.png",
            mime_type="image/png",
            size_bytes=len(b"avatar-bytes"),
        )
    )
    base_video_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="avatar_base_video",
        storage_key=base_video_key,
        filename="avatar-base.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"base-video-bytes"),
        duration_seconds=8.0,
        metadata_json={
            "avatar_source_storage_key": avatar_key,
            "base_video_provider": "wavespeed_infinitetalk_fast",
            "resolution": "480p",
            "base_video_duration_seconds": 8.0,
            "model_used": "wavespeed-ai/sync-lipsync-3",
            "generated_from": "real_avatar_base_video",
            "fallback_used": False,
            "avatar_source_signature": "signature",
        },
    )
    db_session.add(base_video_asset)
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "generation_job_id": "job-123",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "text": "Narration text for fallback",
                    "word_count": 4,
                    "estimated_duration_seconds": 6.0,
                    "measured_tts_duration": 6.0,
                    "audio_storage_key": audio_key,
                    "audio_url": "https://wavespeed.test/uploaded-audio.mp3",
                }
            ],
        },
    )
    db_session.add(audio_asset)
    db_session.commit()

    class FakeProvider:
        def generate_avatar_video_from_base_video(self, **_kwargs):
            raise AvatarVideoProviderError("WaveSpeed prediction timed out", "WAVESPEED_AVATAR_FAILED")

    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_GENERATION_MODE", "fast_lipsync")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_LIPSYNC_PROVIDER", "wavespeed_sync_lipsync_3")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.ENABLE_STATIC_AVATAR_FALLBACK", True)
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_PROVIDER_MAX_RETRIES", 1)
    monkeypatch.setattr("app.modules.generation.pipeline.get_avatar_video_provider", lambda *_args, **_kwargs: FakeProvider())
    monkeypatch.setattr(
        "app.modules.generation.pipeline._ensure_avatar_base_video_asset",
        lambda **_kwargs: (base_video_asset, b"base-video-bytes", base_video_asset.metadata_json or {}),
    )
    monkeypatch.setattr("app.modules.generation.pipeline._image_to_video_clip", lambda image_bytes, duration_seconds: b"base-video-bytes")
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: (
        {
            "duration_seconds": 6.0,
            "has_video": True,
            "has_audio": False,
            "video_codec": "h264",
            "audio_codec": None,
            "width": 1280,
            "height": 720,
        }
        if media_bytes in {b"base-video-bytes", b"static-video-bytes"}
        else {
            "duration_seconds": 6.0,
            "has_video": False,
            "has_audio": True,
            "video_codec": None,
            "audio_codec": "aac",
            "width": None,
            "height": None,
        }
    ))
    monkeypatch.setattr("app.modules.generation.pipeline._static_avatar_fallback_clip", lambda *_args, **_kwargs: (b"static-video-bytes", "provider timeout"))
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.get", lambda url, timeout=None: type("R", (), {"content": b"static-video-bytes", "status_code": 200})())
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.head", lambda url, timeout=None, follow_redirects=True: type("R", (), {"status_code": 200})())
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.strip_audio_from_video", lambda self, video_bytes: video_bytes)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_avatar_clip_for_slide(
        db_session,
        storage,
        job,
        context,
        slide,
        1,
        1,
        audio_asset,
    )

    assert asset.metadata_json["fallback_used"] is True
    assert asset.metadata_json["fallback_reason"] == "provider timeout"
    assert asset.metadata_json["avatar_overlay_type"] == "static_avatar_fallback"
    assert asset.metadata_json["provider_lipsync_output_storage_key"] is None
    assert asset.metadata_json["final_avatar_output_storage_key"] == asset.storage_key


def test_generate_avatar_clip_retry_success_is_not_marked_as_fallback(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")
    base_video_key = f"projects/{uuid.uuid4()}/avatar/avatar-base.mp4"
    storage.upload_file(base_video_key, b"base-video-bytes", "video/mp4")
    audio_key = f"projects/{uuid.uuid4()}/audio/slide-1.mp3"
    storage.upload_file(audio_key, b"audio-bytes", "audio/mpeg")
    slide_preview_key = f"presentations/{uuid.uuid4()}/previews/slide-1.png"
    storage.upload_file(slide_preview_key, b"slide-preview", "image/png")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Fast Lipsync Retry Success Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narration text",
        thumbnail_key=slide_preview_key,
        metadata_={
            "dialogue": "Narration text for retry success",
            "rendered_image_key": slide_preview_key,
            "slide_preview": {
                "asset_type": "slide_preview",
                "storage_key": slide_preview_key,
                "render_source": "ppt_render",
                "includes_text": True,
            },
        },
    )
    db_session.add(slide)
    db_session.flush()
    db_session.add(
        Asset(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            slide_id=None,
            asset_type="avatar_source",
            storage_key=avatar_key,
            filename="avatar.png",
            mime_type="image/png",
            size_bytes=len(b"avatar-bytes"),
        )
    )
    base_video_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="avatar_base_video",
        storage_key=base_video_key,
        filename="avatar-base.mp4",
        mime_type="video/mp4",
        size_bytes=len(b"base-video-bytes"),
        duration_seconds=8.0,
        metadata_json={
            "avatar_source_storage_key": avatar_key,
            "base_video_provider": "wavespeed_infinitetalk_fast",
            "resolution": "480p",
            "base_video_duration_seconds": 8.0,
            "model_used": "wavespeed-ai/sync-lipsync-3",
            "generated_from": "real_avatar_base_video",
            "fallback_used": False,
            "avatar_source_signature": "signature",
        },
    )
    db_session.add(base_video_asset)
    audio_asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key=audio_key,
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=len(b"audio-bytes"),
        duration_seconds=6.0,
        metadata_json={
            "generation_job_id": "job-123",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "text": "Narration text for retry success",
                    "word_count": 4,
                    "estimated_duration_seconds": 6.0,
                    "measured_tts_duration": 6.0,
                    "audio_storage_key": audio_key,
                    "audio_url": "https://wavespeed.test/uploaded-audio.mp3",
                }
            ],
        },
    )
    db_session.add(audio_asset)
    db_session.commit()

    class RetryProvider:
        def __init__(self):
            self.calls = 0
            self.last_request_id = None
            self.last_image_url = None
            self.last_audio_url = None
            self.last_external_checks = {}
            self.last_duration_ratio = None

        def generate_avatar_video_from_base_video(self, **_kwargs):
            self.calls += 1
            if self.calls == 1:
                raise AvatarVideoProviderError("WaveSpeed prediction timed out", "WAVESPEED_AVATAR_FAILED")
            self.last_request_id = "request-retry-success"
            self.last_image_url = "https://wavespeed.test/uploaded-image.mp4"
            self.last_audio_url = "https://wavespeed.test/uploaded-audio.mp3"
            self.last_external_checks = {
                "image": {"validated": True, "status_code": 200},
                "audio": {"validated": True, "status_code": 200},
            }
            self.last_duration_ratio = 1.0
            return "https://cdn.test/retry-success.mp4"

    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_GENERATION_MODE", "fast_lipsync")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_LIPSYNC_PROVIDER", "wavespeed_sync_lipsync_3")
    monkeypatch.setattr("app.modules.generation.pipeline.settings.ENABLE_STATIC_AVATAR_FALLBACK", True)
    monkeypatch.setattr("app.modules.generation.pipeline.settings.AVATAR_PROVIDER_MAX_RETRIES", 1)
    monkeypatch.setattr("app.modules.generation.pipeline.get_avatar_video_provider", lambda *_args, **_kwargs: RetryProvider())
    monkeypatch.setattr("app.modules.generation.pipeline._ensure_avatar_base_video_asset", lambda **_kwargs: (base_video_asset, b"base-video-bytes", base_video_asset.metadata_json or {}))
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.get", lambda url, timeout=None: type("R", (), {"content": b"retry-success-video", "status_code": 200})())
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.head", lambda url, timeout=None, follow_redirects=True: type("R", (), {"status_code": 200})())
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 6.0,
        "has_video": True,
        "has_audio": False,
        "video_stream_count": 1,
        "audio_stream_count": 0,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    } if media_bytes in {b"base-video-bytes", b"retry-success-video"} else {
        "duration_seconds": 6.0,
        "has_video": False,
        "has_audio": True,
        "video_stream_count": 0,
        "audio_stream_count": 1,
        "video_codec": None,
        "audio_codec": "aac",
        "width": None,
        "height": None,
    })
    monkeypatch.setattr("app.modules.generation.pipeline._static_avatar_fallback_clip", lambda *_args, **_kwargs: (_ for _ in ()).throw(AssertionError("static fallback should not be used on retry success")))
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.strip_audio_from_video", lambda self, video_bytes: video_bytes)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_avatar_clip_for_slide(
        db_session,
        storage,
        job,
        context,
        slide,
        1,
        1,
        audio_asset,
        avatar_base_video_asset=base_video_asset,
        avatar_base_video_bytes=b"base-video-bytes",
        avatar_base_video_metadata=base_video_asset.metadata_json or {},
    )

    assert asset.metadata_json["fallback_used"] is False
    assert asset.metadata_json["provider_lipsync_output_storage_key"] == asset.storage_key
    assert asset.metadata_json["final_avatar_output_storage_key"] == asset.storage_key
    assert asset.metadata_json["avatar_overlay_type"] == "provider_lipsync_video"
    assert asset.metadata_json["provider_prediction_id"] == "request-retry-success"
    assert asset.metadata_json["provider_lipsync_output_present"] is True


def test_avatar_base_video_asset_is_reused(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Base Video Reuse Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=uuid.uuid4(),
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{uuid.uuid4()}",
    )

    monkeypatch.setattr("app.modules.generation.pipeline._image_to_video_clip", lambda image_bytes, duration_seconds: b"base-video-bytes")
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", lambda media_bytes, *_args, **_kwargs: {
        "duration_seconds": 12.0,
        "has_video": True,
        "has_audio": False,
        "video_codec": "h264",
        "audio_codec": None,
        "width": 1280,
        "height": 720,
    })
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    asset1, bytes1, _meta1 = _ensure_avatar_base_video_asset(db_session, storage, context)
    asset2, bytes2, _meta2 = _ensure_avatar_base_video_asset(db_session, storage, context)

    assert asset1.id == asset2.id
    assert bytes1 == bytes2 == b"base-video-bytes"
    assert asset1.asset_type == "avatar_base_video"
    assert asset1.metadata_json["avatar_source_storage_key"] == avatar_key


def test_avatar_base_video_real_generation_is_cached_and_reused(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Real Base Video Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=uuid.uuid4(),
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{uuid.uuid4()}",
    )

    class FakeTTSProvider:
        def generate_audio(self, **_kwargs):
            return b"calibration-audio", 5.0

    class FakeWaveSpeedProvider:
        def generate_avatar_video_from_audio(self, **_kwargs):
            self.last_image_url = "https://wavespeed.test/uploaded-image.png"
            self.last_audio_url = "https://wavespeed.test/uploaded-audio.mp3"
            self.last_request_id = "request-real-base"
            self.last_external_checks = {
                "image": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/uploaded-image.png"},
                "audio": {"validated": True, "status_code": 200, "host": "wavespeed.test", "path": "/uploaded-audio.mp3"},
            }
            self.last_duration_ratio = 1.0
            self.last_avatar_duration_seconds = 8.0
            self.last_audio_duration_seconds = 5.0
            return "https://cdn.test/base.mp4"

    def fake_get(url, timeout=None):
        if url == "https://cdn.test/base.mp4":
            return type("R", (), {"status_code": 200, "content": b"base-video-bytes"})()
        raise AssertionError(f"unexpected URL: {url}")

    def fake_probe(media_bytes, *_args, **_kwargs):
        if media_bytes == b"calibration-audio":
            return {
                "duration_seconds": 5.0,
                "has_video": False,
                "has_audio": True,
                "video_codec": None,
                "audio_codec": "aac",
                "width": None,
                "height": None,
            }
        if media_bytes == b"base-video-bytes":
            return {
                "duration_seconds": 8.0,
                "has_video": True,
                "has_audio": False,
                "video_codec": "h264",
                "audio_codec": None,
                "width": 1280,
                "height": 720,
            }
        raise AssertionError(f"unexpected media bytes: {media_bytes!r}")

    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.TTS_PROVIDER", "elevenlabs")
    monkeypatch.setattr("app.modules.generation.pipeline.get_tts_provider", lambda *_args, **_kwargs: FakeTTSProvider())
    monkeypatch.setattr("app.modules.generation.pipeline.get_avatar_video_provider", lambda *_args, **_kwargs: FakeWaveSpeedProvider())
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.normalize_audio_to_mp3", lambda self, audio_bytes: audio_bytes)
    monkeypatch.setattr("app.modules.generation.pipeline.ComposerService.strip_audio_from_video", lambda self, video_bytes: video_bytes)
    monkeypatch.setattr("app.modules.generation.pipeline.httpx.get", fake_get)
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", fake_probe)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    asset1, bytes1, meta1 = _ensure_avatar_base_video_asset(db_session, storage, context)
    asset2, bytes2, meta2 = _ensure_avatar_base_video_asset(db_session, storage, context)

    assert asset1.id == asset2.id
    assert bytes1 == bytes2 == b"base-video-bytes"
    assert meta1["generated_from"] == "real_avatar_base_video"
    assert meta1["fallback_used"] is False


def test_generate_audio_for_slide_accepts_small_overflow_with_tolerance(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Audio Tolerance Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes="Narración breve de prueba.",
        metadata_={"dialogue": "Narración breve de prueba."},
    )
    db_session.add(slide)
    db_session.commit()

    class FakeProvider:
        def generate_audio(self, **_kwargs):
            return b"audio-bytes", 0.0

    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.TTS_PROVIDER", "elevenlabs")
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_LIPSYNC_AUDIO_SECONDS_PER_CHUNK", 20)
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_AUDIO_CHUNK_DURATION_TOLERANCE_SECONDS", 1.0)
    monkeypatch.setattr("app.modules.generation.pipeline.get_tts_provider", lambda *_args, **_kwargs: FakeProvider())
    monkeypatch.setattr(
        "app.modules.generation.pipeline.ComposerService.normalize_audio_to_mp3",
        lambda self, audio_bytes: audio_bytes,
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline._probe_media_info",
        lambda media_bytes, *_args, **_kwargs: {
            "duration_seconds": 20.04,
            "has_video": False,
            "has_audio": True,
            "video_codec": None,
            "audio_codec": "aac",
            "width": None,
            "height": None,
        },
    )
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url="https://example.test/avatar.png",
        avatar_source_storage_key=None,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_audio_for_slide(db_session, storage, job, context, slide, 1, 1)

    assert asset.asset_type == "slide_audio"
    assert asset.duration_seconds == 20.04
    assert asset.metadata_json["chunk_count"] == 1


def test_generate_audio_for_slide_splits_when_chunk_exceeds_tolerance(monkeypatch, db_session):
    storage = InMemoryStorageProvider()
    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Audio Split Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    dialogue = " ".join(["palabra"] * 90)
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes=dialogue,
        metadata_={"dialogue": dialogue},
    )
    db_session.add(slide)
    db_session.commit()

    call_count = {"count": 0}

    class FakeProvider:
        def generate_audio(self, text, **_kwargs):
            call_count["count"] += 1
            return f"audio-{call_count['count']}".encode(), 0.0

    def fake_probe(media_bytes, *_args, **_kwargs):
        mapping = {
            b"audio-1": {
                "duration_seconds": 22.5,
                "has_video": False,
                "has_audio": True,
                "video_codec": None,
                "audio_codec": "aac",
                "width": None,
                "height": None,
            },
            b"audio-2": {
                "duration_seconds": 10.0,
                "has_video": False,
                "has_audio": True,
                "video_codec": None,
                "audio_codec": "aac",
                "width": None,
                "height": None,
            },
            b"audio-3": {
                "duration_seconds": 10.0,
                "has_video": False,
                "has_audio": True,
                "video_codec": None,
                "audio_codec": "aac",
                "width": None,
                "height": None,
            },
            b"joined-audio": {
                "duration_seconds": 20.0,
                "has_video": False,
                "has_audio": True,
                "video_codec": None,
                "audio_codec": "aac",
                "width": None,
                "height": None,
            },
        }
        return mapping[media_bytes]

    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.TTS_PROVIDER", "elevenlabs")
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_LIPSYNC_AUDIO_SECONDS_PER_CHUNK", 20)
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_AUDIO_CHUNK_DURATION_TOLERANCE_SECONDS", 1.0)
    monkeypatch.setattr("app.modules.generation.pipeline.get_tts_provider", lambda *_args, **_kwargs: FakeProvider())
    monkeypatch.setattr(
        "app.modules.generation.pipeline.ComposerService.normalize_audio_to_mp3",
        lambda self, audio_bytes: audio_bytes,
    )
    monkeypatch.setattr(
        "app.modules.generation.pipeline.ComposerService.concatenate_audio_tracks",
        lambda self, tracks: b"joined-audio",
    )
    monkeypatch.setattr("app.modules.generation.pipeline._probe_media_info", fake_probe)
    monkeypatch.setattr("app.modules.generation.pipeline.storage_object_exists", lambda *_: True)

    job = GenerationJob(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        status="queued",
        progress_percentage=0.0,
    )
    db_session.add(job)
    db_session.flush()
    context = GenerationContext(
        generation_job_id=job.id,
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=presentation,
        slides=[slide],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key="secret",
        elevenlabs_voice_id="voice-id",
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url="https://example.test/avatar.png",
        avatar_source_storage_key=None,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{job.id}",
    )

    asset = generate_audio_for_slide(db_session, storage, job, context, slide, 1, 1)

    assert call_count["count"] >= 2
    assert asset.asset_type == "slide_audio"
    assert asset.metadata_json["chunk_count"] >= 2


def test_load_avatar_source_url_uses_public_storage_url(db_session):
    storage = InMemoryStorageProvider()
    avatar_key = f"projects/{uuid.uuid4()}/avatar/avatar.png"
    storage.upload_file(avatar_key, b"avatar-bytes", "image/png")

    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Avatar URL Project",
    )
    db_session.add(project)
    db_session.flush()
    asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="avatar_source",
        storage_key=avatar_key,
        filename="avatar.png",
        mime_type="image/png",
        size_bytes=len(b"avatar-bytes"),
    )
    db_session.add(asset)
    db_session.flush()

    context = GenerationContext(
        generation_job_id=uuid.uuid4(),
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        presentation=None,  # type: ignore[arg-type]
        slides=[],
        settings=VideoGenerationSettings(
            organization_id=MOCK_ORG_ID,
            project_id=project.id,
            validation_status="saved",
        ),
        elevenlabs_api_key=None,
        elevenlabs_voice_id=None,
        wavespeed_api_key="wavespeed-secret",
        avatar_source_url=None,
        avatar_source_storage_key=avatar_key,
        output_prefix=f"orgs/{MOCK_ORG_ID}/projects/{project.id}/generation/{uuid.uuid4()}",
    )

    url, metadata = _load_avatar_source_url(db_session, storage, context)

    assert url == f"https://storage.example.test/download/{avatar_key}"
    assert metadata["storage_key"] == avatar_key
    assert metadata["source"] == "storage"


def test_slide_segment_duration_uses_audio_duration_and_pause():
    audio_asset = Asset(duration_seconds=12.0)
    avatar_asset = Asset(duration_seconds=9.0)

    assert _slide_segment_duration_seconds(audio_asset, avatar_asset) == 12.5
    try:
        _slide_segment_duration_seconds(None, avatar_asset)
    except Exception as exc:
        assert "controlled TTS narration track" in str(exc)
    else:
        raise AssertionError("slide duration must require audio")


def test_normalize_tts_text_adds_sentence_punctuation():
    text = "Hola mundo\n\nEsto es una prueba"
    normalized = _normalize_tts_text(text)

    assert normalized == "Hola mundo. Esto es una prueba."


def test_build_narration_chunks_splits_long_text(monkeypatch):
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_TTS_CHARS_PER_CHUNK", 120)
    monkeypatch.setattr(
        "app.modules.generation.pipeline.app_settings.MAX_LIPSYNC_AUDIO_SECONDS_PER_CHUNK",
        20,
    )
    text = (
        "Este es un párrafo largo con varias oraciones. "
        "Necesitamos dividirlo cuidadosamente sin cortar palabras ni perder acentos. "
        "La narración debe mantenerse natural y clara para la voz en español. "
        "Además, cada fragmento debe quedar dentro de un tamaño razonable. "
    ) * 2

    chunks = _build_narration_chunks(text)

    assert len(chunks) > 1
    assert all(chunk["text_length"] <= 120 for chunk in chunks)
    assert all(chunk["expected_duration_seconds"] <= 20 for chunk in chunks)


def test_build_narration_chunks_caps_chunk_count(monkeypatch):
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_TTS_CHARS_PER_CHUNK", 700)
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_LIPSYNC_AUDIO_SECONDS_PER_CHUNK", 30)
    monkeypatch.setattr("app.modules.generation.pipeline.app_settings.MAX_CHUNKS_PER_SLIDE", 4)
    text = ("Este es un texto largo y natural para probar el límite de fragmentos. " * 30).strip()

    chunks = _build_narration_chunks(text)

    assert len(chunks) <= 4


def test_detect_green_background_identifies_green_border(monkeypatch):
    # 64x64 RGB frame, border pixels are green, center is neutral.
    width = height = 64
    green = bytes((20, 220, 20))
    neutral = bytes((80, 80, 80))
    frame = bytearray()
    for y in range(height):
        for x in range(width):
            if x < 4 or x >= width - 4 or y < 4 or y >= height - 4:
                frame.extend(green)
            else:
                frame.extend(neutral)

    class FakeResult:
        stdout = bytes(frame)

    def fake_run(*_args, **_kwargs):
        return FakeResult()

    monkeypatch.setattr("app.modules.generation.pipeline.subprocess.run", fake_run)

    info = _detect_green_background(b"video-bytes")

    assert info["detected"] is True
    assert info["green_ratio"] > 0.45


def test_detect_green_background_ignores_neutral_frame(monkeypatch):
    width = height = 64
    neutral = bytes((80, 80, 80))
    frame = neutral * (width * height)

    class FakeResult:
        stdout = bytes(frame)

    def fake_run(*_args, **_kwargs):
        return FakeResult()

    monkeypatch.setattr("app.modules.generation.pipeline.subprocess.run", fake_run)

    info = _detect_green_background(b"video-bytes")

    assert info["detected"] is False
    assert info["green_ratio"] == 0.0


def test_can_reuse_slide_audio_asset_rejects_legacy_single_chunk_audio(db_session):
    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Legacy Audio Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    slide = Slide(
        presentation_id=presentation.id,
        position=1,
        title="Slide 1",
        notes=" ".join(["palabra"] * 320),
        metadata_={"dialogue": " ".join(["palabra"] * 320)},
    )
    db_session.add(slide)
    db_session.flush()
    asset = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=slide.id,
        asset_type="slide_audio",
        storage_key="projects/audio/legacy.mp3",
        filename="legacy.mp3",
        mime_type="audio/mpeg",
        size_bytes=100,
        duration_seconds=20.0,
        metadata_json={"generation_job_id": "job-legacy", "slide_position": 1},
    )
    db_session.add(asset)
    db_session.commit()

    assert _can_reuse_slide_audio_asset(asset, slide, 1) is False


def test_slide_audio_chunk_specs_requires_chunk_metadata_for_long_narration():
    audio_asset = Asset(
        duration_seconds=20.0,
        metadata_json={"generation_job_id": "job-legacy", "slide_position": 1},
    )
    dialogue = " ".join(["palabra"] * 320)

    try:
        _slide_audio_chunk_specs(audio_asset, dialogue)
    except Exception as exc:
        assert "missing chunk metadata" in str(exc)
    else:
        raise AssertionError("long narration must require chunk metadata")


def test_slide_audio_chunk_specs_rejects_missing_text_or_audio_fields(db_session):
    project = Project(
        organization_id=MOCK_ORG_ID,
        owner_id=MOCK_USER_ID,
        name="Bad Chunk Project",
    )
    db_session.add(project)
    db_session.flush()
    presentation = Presentation(
        project_id=project.id,
        organization_id=MOCK_ORG_ID,
        title="deck.pptx",
        original_filename="deck.pptx",
        storage_key="projects/deck.pptx",
        status=PresentationStatus.parsed,
        slide_count=1,
    )
    db_session.add(presentation)
    db_session.flush()
    audio_asset_missing_text = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="slide_audio",
        storage_key="projects/audio/slide-1.mp3",
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=100,
        duration_seconds=10.0,
        metadata_json={
            "generation_job_id": "job-1",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "word_count": 2,
                    "estimated_duration": 1.0,
                    "audio_storage_key": "projects/audio/slide-1.mp3",
                    "audio_url": "https://storage.example.test/download/projects/audio/slide-1.mp3",
                }
            ],
        },
    )

    try:
        _slide_audio_chunk_specs(audio_asset_missing_text, "")
    except Exception as exc:
        assert "missing required narration text" in str(exc)
    else:
        raise AssertionError("malformed chunks must be rejected with a clear error")

    audio_asset_missing_url = Asset(
        organization_id=MOCK_ORG_ID,
        project_id=project.id,
        slide_id=None,
        asset_type="slide_audio",
        storage_key="projects/audio/slide-1.mp3",
        filename="slide-1.mp3",
        mime_type="audio/mpeg",
        size_bytes=100,
        duration_seconds=10.0,
        metadata_json={
            "generation_job_id": "job-1",
            "slide_position": 1,
            "chunks": [
                {
                    "index": 1,
                    "text": "Hola mundo",
                    "word_count": 2,
                    "estimated_duration": 1.0,
                    "audio_storage_key": "projects/audio/slide-1.mp3",
                }
            ],
        },
    )

    try:
        _slide_audio_chunk_specs(audio_asset_missing_url, "Hola mundo")
    except Exception as exc:
        assert "missing required audio URL" in str(exc)
    else:
        raise AssertionError("malformed chunks must be rejected with a clear error")
