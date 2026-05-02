import uuid
from io import BytesIO
from unittest.mock import patch

from fastapi.testclient import TestClient
from pptx import Presentation as PptxPresentation

from app.modules.organizations.models import Organization
from app.modules.projects.models import Presentation, PresentationStatus, Project, Slide
from app.modules.projects.service import MOCK_ORG_ID, MOCK_USER_ID
from app.modules.users.models import User
from tests.conftest import _TestingSession
from tests.fakes import InMemoryStorageProvider

BASE = "/api/v1"


def _create_presentation_with_slides() -> tuple[uuid.UUID, list[uuid.UUID]]:
    db = _TestingSession()
    try:
        user = User(
            id=MOCK_USER_ID,
            email=f"{uuid.uuid4()}@test.local",
            hashed_password="test",
            full_name="Test User",
        )
        org = Organization(id=MOCK_ORG_ID, name="Test Org", slug=f"test-{uuid.uuid4()}")
        project = Project(
            organization_id=MOCK_ORG_ID,
            owner_id=MOCK_USER_ID,
            name="Slides Project",
        )
        db.add_all([user, org, project])
        db.flush()

        presentation = Presentation(
            project_id=project.id,
            organization_id=MOCK_ORG_ID,
            title="deck.pptx",
            original_filename="deck.pptx",
            storage_key="test/deck.pptx",
            status=PresentationStatus.parsed,
            slide_count=2,
        )
        db.add(presentation)
        db.flush()

        second = Slide(
            presentation_id=presentation.id,
            position=2,
            title="Second",
            notes="Second notes",
            metadata_={
                "visible_text": "Second visible text",
                "dialogue": "Second dialogue",
                "keep": "yes",
            },
        )
        first = Slide(
            presentation_id=presentation.id,
            position=1,
            title="First",
            notes="First notes",
            metadata_={
                "visible_text": "First visible text",
                "dialogue": "First dialogue",
                "custom": {"level": 1},
            },
        )
        db.add_all([second, first])
        db.commit()
        return presentation.id, [first.id, second.id]
    finally:
        db.close()


def test_list_slides_by_presentation_returns_position_order(client: TestClient):
    presentation_id, slide_ids = _create_presentation_with_slides()

    res = client.get(f"{BASE}/presentations/{presentation_id}/slides")

    assert res.status_code == 200
    body = res.json()
    assert [slide["id"] for slide in body] == [str(slide_ids[0]), str(slide_ids[1])]
    assert [slide["position"] for slide in body] == [1, 2]
    assert body[0]["visible_text"] == "First visible text"
    assert body[0]["dialogue"] == "First dialogue"
    assert body[0]["metadata"]["custom"] == {"level": 1}
    assert "created_at" in body[0]
    assert "updated_at" in body[0]


def test_get_slide_returns_editor_fields(client: TestClient):
    _presentation_id, slide_ids = _create_presentation_with_slides()

    res = client.get(f"{BASE}/slides/{slide_ids[0]}")

    assert res.status_code == 200
    body = res.json()
    assert body["id"] == str(slide_ids[0])
    assert body["title"] == "First"
    assert body["notes"] == "First notes"
    assert body["visible_text"] == "First visible text"
    assert body["dialogue"] == "First dialogue"


def test_patch_slide_updates_dialogue_in_metadata(client: TestClient):
    _presentation_id, slide_ids = _create_presentation_with_slides()

    res = client.patch(
        f"{BASE}/slides/{slide_ids[0]}",
        json={"dialogue": "Updated dialogue"},
    )

    assert res.status_code == 200
    body = res.json()
    assert body["dialogue"] == "Updated dialogue"
    assert body["metadata"]["dialogue"] == "Updated dialogue"
    assert body["metadata"]["visible_text"] == "First visible text"


def test_patch_slide_preserves_existing_metadata(client: TestClient):
    _presentation_id, slide_ids = _create_presentation_with_slides()

    res = client.patch(
        f"{BASE}/slides/{slide_ids[0]}",
        json={
            "title": "Updated title",
            "notes": "Updated notes",
            "visible_text": "Updated visible text",
            "metadata": {"reviewed": True},
        },
    )

    assert res.status_code == 200
    body = res.json()
    assert body["title"] == "Updated title"
    assert body["notes"] == "Updated notes"
    assert body["visible_text"] == "Updated visible text"
    assert body["metadata"]["visible_text"] == "Updated visible text"
    assert body["metadata"]["dialogue"] == "First dialogue"
    assert body["metadata"]["custom"] == {"level": 1}
    assert body["metadata"]["reviewed"] is True


def test_patch_slide_text_updates_pptx_and_preview_key(client: TestClient):
    storage = InMemoryStorageProvider()
    storage_key = f"{MOCK_ORG_ID}/{uuid.uuid4()}/deck.pptx"
    storage.upload_file(
        storage_key,
        _make_single_slide_pptx("Original title"),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    slide_id = _create_presentation_with_editable_slide(storage_key)

    metadata = {
        "visible_text": "Updated title",
        "dialogue": "Narration stays separate",
        "canvas": {
            "width": 960,
            "height": 540,
            "text": {"title": "Updated title", "visible_text": "Updated title"},
            "text_blocks": [
                {
                    "id": "title-0",
                    "shape_index": 0,
                    "type": "title",
                    "text": "Updated title",
                    "x": 10,
                    "y": 10,
                    "width": 400,
                    "height": 80,
                    "fontSize": 34,
                    "fontWeight": "700",
                    "color": "#111827",
                    "textAlign": "left",
                }
            ],
        },
    }

    with patch("app.modules.slides.service.get_storage", return_value=storage), patch(
        "app.modules.slides.service.render_slide_previews",
        return_value={1: "new-preview-key.png"},
    ):
        res = client.patch(
            f"{BASE}/slides/{slide_id}",
            json={
                "title": "Updated title",
                "visible_text": "Updated title",
                "metadata": metadata,
            },
        )

    assert res.status_code == 200
    body = res.json()
    assert body["thumbnail_key"] == "new-preview-key.png"
    assert body["metadata"]["visible_text"] == "Updated title"
    assert body["metadata"]["canvas"]["text_blocks"][0]["text"] == "Updated title"

    updated_deck = PptxPresentation(BytesIO(storage.download_file(storage_key)))
    assert updated_deck.slides[0].shapes[0].text == "Updated title"


def _make_single_slide_pptx(title: str) -> bytes:
    deck = PptxPresentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = title
    buffer = BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _create_presentation_with_editable_slide(storage_key: str) -> uuid.UUID:
    db = _TestingSession()
    try:
        user = User(
            id=MOCK_USER_ID,
            email=f"{uuid.uuid4()}@test.local",
            hashed_password="test",
            full_name="Test User",
        )
        org = Organization(id=MOCK_ORG_ID, name="Test Org", slug=f"test-{uuid.uuid4()}")
        project = Project(
            organization_id=MOCK_ORG_ID,
            owner_id=MOCK_USER_ID,
            name="Editable Slides Project",
        )
        db.add_all([user, org, project])
        db.flush()

        presentation = Presentation(
            project_id=project.id,
            organization_id=MOCK_ORG_ID,
            title="deck.pptx",
            original_filename="deck.pptx",
            storage_key=storage_key,
            status=PresentationStatus.parsed,
            slide_count=1,
        )
        db.add(presentation)
        db.flush()

        slide = Slide(
            presentation_id=presentation.id,
            position=1,
            title="Original title",
            metadata_={
                "visible_text": "Original title",
                "dialogue": "Narration stays separate",
            },
        )
        db.add(slide)
        db.commit()
        return slide.id
    finally:
        db.close()
