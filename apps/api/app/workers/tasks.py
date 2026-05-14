import io
import ipaddress
import json
import logging
import subprocess
import tempfile
import uuid
import xml.etree.ElementTree as ET
from datetime import UTC, datetime
from pathlib import Path
from urllib.parse import urlparse

from celery.exceptions import SoftTimeLimitExceeded
from pptx import Presentation as PptxPresentation
from sqlalchemy import delete

import app.models  # noqa: F401
from app.config import settings
from app.modules.composer.service import ComposerService
from app.modules.generation.models import GenerationJob, VideoGenerationSettings
from app.modules.generation.pipeline import (
    PipelineError,
    compose_final_video,
    compose_segment_for_slide,
    generate_audio_for_slide,
    generate_avatar_clip_for_slide,
    mark_generation_failed,
    validate_generation_job,
)
from app.modules.presentations.rendering import render_slide_previews
from app.modules.projects.models import Asset, Presentation, PresentationStatus, Slide
from app.modules.tts.adapters import TTSProviderError, get_tts_provider
from app.modules.video.adapters import AvatarVideoProviderError, get_avatar_video_provider
from app.providers.storage import get_storage
from app.utils.crypto import decrypt_secret
from app.workers.base_task import JobTask
from app.workers.celery_app import celery_app
from app.workers.db import worker_db_session

logger = logging.getLogger(__name__)

DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATIONML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
XML_NS = {"a": DRAWINGML_NS, "p": PRESENTATIONML_NS}
SCHEME_COLOR_ALIASES = {
    "tx1": "dk1",
    "tx2": "dk2",
    "bg1": "lt1",
    "bg2": "lt2",
    "text1": "dk1",
    "text2": "dk2",
    "background1": "lt1",
    "background2": "lt2",
    "accent_1": "accent1",
    "accent_2": "accent2",
    "accent_3": "accent3",
    "accent_4": "accent4",
    "accent_5": "accent5",
    "accent_6": "accent6",
}


# ── ping ──────────────────────────────────────────────────────────────────────

@celery_app.task(name="app.workers.tasks.ping", bind=True)
def ping(self, message: str = "pong") -> dict:
    """Task de prueba — verifica que el worker está vivo y puede procesar tasks."""
    logger.info("ping recibido: %s", message)
    return {"message": message, "worker": self.request.hostname}


# ── parse_presentation ────────────────────────────────────────────────────────

class ParsePresentationTask(JobTask):
    """Parsea una presentación PPTX y crea registros Slide."""

    name = "app.workers.tasks.parse_presentation"

    def run_job(self, job_id: uuid.UUID, presentation_id: str, **kwargs) -> dict:
        presentation_uuid = uuid.UUID(presentation_id)
        logger.info("Iniciando parse de presentación %s", presentation_id)

        try:
            self.set_progress(job_id, 10.0, "Loading presentation record")
            with worker_db_session() as db:
                presentation = db.get(Presentation, presentation_uuid)
                if not presentation:
                    raise ValueError(f"Presentation not found: {presentation_id}")

                presentation.status = PresentationStatus.processing
                storage_key = presentation.storage_key
                original_filename = presentation.original_filename
                db.commit()

            self.set_progress(job_id, 25.0, "Downloading presentation from storage")
            storage = get_storage()
            pptx_bytes = storage.download_file(storage_key)

            self.set_progress(job_id, 45.0, "Opening PPTX")
            deck = PptxPresentation(io.BytesIO(pptx_bytes))
            theme_colors = _extract_theme_colors(deck)

            self.set_progress(job_id, 65.0, "Extracting slides")
            slide_records = [
                _extract_slide(
                    slide=slide,
                    slide_number=index,
                    slide_width=int(deck.slide_width),
                    slide_height=int(deck.slide_height),
                    theme_colors=theme_colors,
                )
                for index, slide in enumerate(deck.slides, 1)
            ]

            self.set_progress(job_id, 75.0, "Rendering slide previews")
            preview_keys = render_slide_previews(
                pptx_bytes=pptx_bytes,
                presentation_id=presentation_uuid,
                original_filename=original_filename,
                storage=storage,
            )

            self.set_progress(job_id, 85.0, "Saving parsed slides")
            with worker_db_session() as db:
                presentation = db.get(Presentation, presentation_uuid)
                if not presentation:
                    raise ValueError(f"Presentation not found: {presentation_id}")

                db.execute(delete(Slide).where(Slide.presentation_id == presentation_uuid))
                for record in slide_records:
                    db.add(
                        Slide(
                            presentation_id=presentation_uuid,
                            position=record["slide_number"],
                            title=record["title"],
                            notes=record["speaker_notes"] or None,
                            thumbnail_key=preview_keys.get(record["slide_number"]),
                            metadata_={
                                "slide_number": record["slide_number"],
                                "visible_text": record["visible_text"],
                                "dialogue": record["dialogue"],
                                "rendered_image_key": preview_keys.get(record["slide_number"]),
                                "slide_preview": {
                                    "asset_type": "slide_preview",
                                    "storage_key": preview_keys.get(record["slide_number"]),
                                    "render_source": "ppt_render",
                                    "includes_text": True,
                                },
                            },
                        )
                    )

                presentation.slide_count = len(slide_records)
                presentation.status = PresentationStatus.parsed
                db.commit()

            self.set_progress(job_id, 100.0, "Parsed")
            return {
                "presentation_id": presentation_id,
                "slide_count": len(slide_records),
                "parsed": True,
            }
        except Exception as exc:
            error_message = str(exc)
            with worker_db_session() as db:
                presentation = db.get(Presentation, presentation_uuid)
                if presentation:
                    presentation.status = PresentationStatus.failed
                    db.commit()
            logger.exception("Parse de presentación %s falló: %s", presentation_id, exc)
            raise RuntimeError(error_message) from exc


def _extract_slide(
    slide,
    slide_number: int,
    slide_width: int,
    slide_height: int,
    theme_colors: dict[str, str] | None = None,
) -> dict:
    visible_text = _extract_visible_text(slide)
    speaker_notes = _extract_speaker_notes(slide)
    title = _extract_title(slide, visible_text)
    text_blocks = _extract_text_blocks(slide, slide_width, slide_height, theme_colors or {})
    elements = _text_blocks_to_elements(text_blocks)

    return {
        "slide_number": slide_number,
        "title": title,
        "visible_text": visible_text,
        "speaker_notes": speaker_notes,
        "dialogue": speaker_notes,
        "text_blocks": text_blocks,
        "elements": elements,
    }


def _extract_visible_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _normalize_text(shape.text)
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_speaker_notes(slide) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""

    notes_slide = slide.notes_slide
    notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
    if notes_text_frame is None:
        return ""
    return _normalize_text(notes_text_frame.text)


def _extract_title(slide, visible_text: str) -> str | None:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title = _normalize_text(title_shape.text)
        if title:
            return title[:500]

    for line in visible_text.splitlines():
        if line.strip():
            return line.strip()[:500]
    return None


def _extract_text_blocks(
    slide,
    slide_width: int,
    slide_height: int,
    theme_colors: dict[str, str] | None = None,
) -> list[dict]:
    blocks: list[dict] = []
    title_shape = getattr(slide.shapes, "title", None)
    canvas_width = 960
    canvas_height = _scale_emu(slide_height, slide_width, canvas_width) or 540

    for index, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue

        text = _normalize_text(shape.text)
        if not text:
            continue

        block_type = "title" if title_shape is not None and shape == title_shape else "body"
        style = _extract_text_style(shape, block_type, theme_colors or {})
        blocks.append(
            {
                "id": f"{block_type}-{index}",
                "type": block_type,
                "text": text,
                "shape_index": index,
                "x": _scale_emu(getattr(shape, "left", 0), slide_width, canvas_width),
                "y": _scale_emu(getattr(shape, "top", 0), slide_height, canvas_height),
                "width": _scale_emu(getattr(shape, "width", 0), slide_width, canvas_width),
                "height": _scale_emu(getattr(shape, "height", 0), slide_height, canvas_height),
                "fontSize": style["fontSize"],
                "fontWeight": style["fontWeight"],
                "fontFamily": style["fontFamily"],
                "originalFontFamily": style["originalFontFamily"],
                "fallbackFontFamily": style["fallbackFontFamily"],
                "color": style["color"],
                "originalColor": style["originalColor"],
                "textAlign": style["textAlign"],
                "lineHeight": style["lineHeight"],
                "letterSpacing": style["letterSpacing"],
                "bold": style["bold"],
                "italic": style["italic"],
                "underline": style["underline"],
                "style": style,
            }
        )

    return blocks


def _text_blocks_to_elements(text_blocks: list[dict]) -> list[dict]:
    elements: list[dict] = []
    for z_index, block in enumerate(text_blocks, 10):
        elements.append(
            {
                "id": f"text-{block['shape_index']}",
                "type": "text",
                "role": block["type"],
                "shape_index": block["shape_index"],
                "x": block["x"],
                "y": block["y"],
                "width": block["width"],
                "height": block["height"],
                "rotation": 0,
                "zIndex": z_index,
                "text": block["text"],
                "style": block["style"],
            }
        )
    return elements


def _extract_text_style(shape, block_type: str, theme_colors: dict[str, str]) -> dict:
    font = None
    paragraph = None
    run = None
    if getattr(shape, "has_text_frame", False):
        for candidate in shape.text_frame.paragraphs:
            if _normalize_text(candidate.text):
                paragraph = candidate
                for candidate_run in candidate.runs:
                    if _normalize_text(candidate_run.text):
                        run = candidate_run
                        font = candidate_run.font
                        break
                if font is None:
                    font = getattr(candidate, "font", None)
                break

    original_font = _font_name(font)
    bold = bool(getattr(font, "bold", False)) if font else block_type == "title"
    italic = bool(getattr(font, "italic", False)) if font else False
    underline = bool(getattr(font, "underline", False)) if font else False
    font_size = _font_size(font, 34 if block_type == "title" else 22)
    font_weight = "700" if bold or block_type == "title" else "400"
    color = _resolve_text_color(
        run=run,
        paragraph=paragraph,
        shape=shape,
        font=font,
        theme_colors=theme_colors,
        fallback="#111827" if block_type == "title" else "#1f2937",
    )
    return {
        "fontFamily": original_font or "Arial",
        "originalFontFamily": original_font,
        "fallbackFontFamily": "Arial",
        "fontSize": font_size,
        "fontWeight": font_weight,
        "color": color,
        "originalColor": color,
        "bold": bold,
        "italic": italic,
        "underline": underline,
        "textAlign": _paragraph_alignment(paragraph),
        "lineHeight": 1.15,
        "letterSpacing": 0,
        "backgroundColor": "transparent",
    }


def _font_name(font) -> str | None:
    name = getattr(font, "name", None) if font else None
    return str(name) if name else None


def _font_size(font, fallback: int) -> int:
    size = getattr(font, "size", None) if font else None
    if size is None:
        return fallback
    return max(8, round(float(size.pt)))


def _resolve_text_color(
    run,
    paragraph,
    shape,
    font,
    theme_colors: dict[str, str],
    fallback: str,
) -> str:
    color = _color_from_run_xml(run, theme_colors)
    if color:
        return color

    color = _font_color(font, theme_colors)
    if color:
        return color

    color = _color_from_paragraph_xml(paragraph, theme_colors)
    if color:
        return color

    color = _color_from_shape_xml(shape, theme_colors)
    if color:
        return color

    color = _color_from_placeholder_styles(shape, theme_colors)
    if color:
        return color

    return fallback


def _font_color(font, theme_colors: dict[str, str]) -> str | None:
    color = getattr(font, "color", None) if font else None
    try:
        rgb = color.rgb if color else None
    except AttributeError:
        rgb = None
    if rgb:
        return f"#{rgb}".upper()

    try:
        theme_color = color.theme_color if color else None
    except AttributeError:
        theme_color = None
    if theme_color:
        return _resolve_scheme_color(str(getattr(theme_color, "name", theme_color)), theme_colors)
    return None


def _color_from_run_xml(run, theme_colors: dict[str, str]) -> str | None:
    if run is None:
        return None
    run_element = getattr(run, "_r", None)
    if run_element is None:
        return None
    rpr = run_element.find("a:rPr", namespaces=XML_NS)
    return _color_from_solid_fill(rpr, theme_colors)


def _color_from_paragraph_xml(paragraph, theme_colors: dict[str, str]) -> str | None:
    if paragraph is None:
        return None
    paragraph_element = getattr(paragraph, "_p", None)
    if paragraph_element is None:
        return None
    for path in ("a:pPr/a:defRPr", "a:pPr/a:endParaRPr"):
        color = _color_from_solid_fill(
            paragraph_element.find(path, namespaces=XML_NS),
            theme_colors,
        )
        if color:
            return color
    return None


def _color_from_shape_xml(shape, theme_colors: dict[str, str]) -> str | None:
    shape_element = getattr(shape, "_element", None)
    if shape_element is None:
        return None
    for path in (
        ".//a:lstStyle/a:lvl1pPr/a:defRPr",
        ".//a:lstStyle/a:lvl2pPr/a:defRPr",
        ".//a:lstStyle/a:defPPr/a:defRPr",
        ".//a:defRPr",
    ):
        color = _color_from_solid_fill(shape_element.find(path, namespaces=XML_NS), theme_colors)
        if color:
            return color
    return None


def _color_from_placeholder_styles(shape, theme_colors: dict[str, str]) -> str | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        placeholder_idx = shape.placeholder_format.idx
        layout = shape.part.slide_layout
    except (AttributeError, KeyError):
        return None

    for placeholder_collection in (
        getattr(layout, "placeholders", []),
        getattr(getattr(layout, "slide_master", None), "placeholders", []),
    ):
        for placeholder in placeholder_collection:
            try:
                if placeholder.placeholder_format.idx != placeholder_idx:
                    continue
            except (AttributeError, KeyError):
                continue
            color = _color_from_shape_xml(placeholder, theme_colors)
            if color:
                return color
    return None


def _color_from_solid_fill(element, theme_colors: dict[str, str]) -> str | None:
    if element is None:
        return None
    solid_fill = element.find(".//a:solidFill", namespaces=XML_NS)
    if solid_fill is None and _local_name(getattr(element, "tag", "")) == "solidFill":
        solid_fill = element
    if solid_fill is None:
        return None

    srgb = solid_fill.find("a:srgbClr", namespaces=XML_NS)
    if srgb is not None:
        value = srgb.get("val")
        if value:
            return _apply_luminance(_hex(value), srgb)

    scheme = solid_fill.find("a:schemeClr", namespaces=XML_NS)
    if scheme is not None:
        value = scheme.get("val")
        resolved = _resolve_scheme_color(value, theme_colors)
        if resolved:
            return _apply_luminance(resolved, scheme)

    return None


def _extract_theme_colors(deck) -> dict[str, str]:
    colors: dict[str, str] = {}
    try:
        theme_parts = [
            part
            for part in deck.part.package.iter_parts()
            if "theme" in str(part.partname)
        ]
    except AttributeError:
        return colors

    for theme_part in theme_parts:
        try:
            root = ET.fromstring(theme_part.blob)
        except ET.ParseError:
            continue
        color_scheme = root.find(".//a:clrScheme", namespaces=XML_NS)
        if color_scheme is None:
            continue
        for child in list(color_scheme):
            name = _local_name(child.tag)
            srgb = child.find("a:srgbClr", namespaces=XML_NS)
            if srgb is not None and srgb.get("val"):
                colors[name] = _hex(srgb.get("val"))
                continue
            sys_color = child.find("a:sysClr", namespaces=XML_NS)
            if sys_color is not None and sys_color.get("lastClr"):
                colors[name] = _hex(sys_color.get("lastClr"))
    return colors


def _resolve_scheme_color(value: str | None, theme_colors: dict[str, str]) -> str | None:
    if not value:
        return None
    key = str(value)
    normalized = key.lower().replace(" ", "").replace("-", "_")
    candidates = [
        key,
        normalized,
        normalized.replace("_", ""),
        SCHEME_COLOR_ALIASES.get(normalized),
        SCHEME_COLOR_ALIASES.get(normalized.replace("_", "")),
    ]
    for candidate in candidates:
        if candidate and candidate in theme_colors:
            return theme_colors[candidate]
    return None


def _apply_luminance(hex_color: str, element) -> str:
    rgb = _rgb_tuple(hex_color)
    if rgb is None:
        return hex_color

    lum_mod = _color_modifier(element, "lumMod", 100000)
    lum_off = _color_modifier(element, "lumOff", 0)
    adjusted = [
        round(channel * (lum_mod / 100000) + 255 * (lum_off / 100000))
        for channel in rgb
    ]
    return _rgb_hex(tuple(clamp for clamp in (_clamp_channel(value) for value in adjusted)))


def _color_modifier(element, name: str, fallback: int) -> int:
    child = element.find(f"a:{name}", namespaces=XML_NS)
    if child is None:
        return fallback
    try:
        return int(child.get("val", fallback))
    except (TypeError, ValueError):
        return fallback


def _hex(value: str | None) -> str:
    if not value:
        return "#000000"
    clean = value.strip().lstrip("#").upper()
    if len(clean) == 3:
        clean = "".join(char * 2 for char in clean)
    return f"#{clean[:6].ljust(6, '0')}"


def _rgb_tuple(hex_color: str) -> tuple[int, int, int] | None:
    clean = hex_color.strip().lstrip("#")
    if len(clean) != 6:
        return None
    try:
        return tuple(int(clean[index : index + 2], 16) for index in (0, 2, 4))
    except ValueError:
        return None


def _rgb_hex(rgb: tuple[int, int, int]) -> str:
    return "#{:02X}{:02X}{:02X}".format(*rgb)


def _clamp_channel(value: int) -> int:
    return max(0, min(255, value))


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _paragraph_alignment(paragraph) -> str:
    alignment = getattr(paragraph, "alignment", None)
    if alignment is None:
        return "left"
    name = getattr(alignment, "name", "")
    if name == "CENTER":
        return "center"
    if name == "RIGHT":
        return "right"
    if name in {"JUSTIFY", "DISTRIBUTE"}:
        return "justify"
    return "left"


def _scale_emu(value: int, source_size: int, target_size: int) -> int:
    if source_size <= 0:
        return 0
    return round((int(value) / source_size) * target_size)


def _normalize_text(text: str | None) -> str:
    if not text:
        return ""
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


# Registrar la task en Celery
parse_presentation = celery_app.register_task(ParsePresentationTask())


class GenerateVideoTask(JobTask):
    name = "app.workers.tasks.generate_video"
    soft_time_limit = 1800
    time_limit = 2100

    def run_job(
        self,
        job_id: uuid.UUID,
        generation_job_id: str,
        project_id: str,
        **kwargs,
    ) -> dict:
        generation_uuid = uuid.UUID(generation_job_id)
        project_uuid = uuid.UUID(project_id)
        logger.info(
            "generate_video task started: project_id=%s generation_job_id=%s celery_request_id=%s",
            project_uuid,
            generation_uuid,
            self.request.id,
        )
        storage = get_storage()
        _log_storage_config_for_generation()

        try:
            with worker_db_session() as db:
                generation_job = db.get(GenerationJob, generation_uuid)
                if not generation_job:
                    raise RuntimeError("Generation job not found")
                context = validate_generation_job(
                    db=db,
                    job=generation_job,
                    project_id=project_uuid,
                    organization_id=generation_job.organization_id,
                    storage=storage,
                )

                total_slides = len(context.slides)
                avatar_clip_assets = [
                    generate_avatar_clip_for_slide(
                        db,
                        storage,
                        generation_job,
                        context,
                        slide,
                        slide_index,
                        total_slides,
                        None,
                    )
                    for slide_index, slide in enumerate(context.slides, 1)
                ]

                segment_assets = [
                    compose_segment_for_slide(
                        db,
                        storage,
                        generation_job,
                        context,
                        slide,
                        slide_index,
                        total_slides,
                        None,
                        avatar_clip_assets[slide_index - 1],
                    )
                    for slide_index, slide in enumerate(context.slides, 1)
                ]

                final_asset = compose_final_video(
                    db,
                    storage,
                    generation_job,
                    context,
                    segment_assets,
                )
                result = {
                    "final_video_key": final_asset.storage_key,
                    "final_asset_id": str(final_asset.id),
                }
            self.set_progress(job_id, 100.0, "Completed")
            return result
        except SoftTimeLimitExceeded as exc:
            with worker_db_session() as db:
                generation_job = db.get(GenerationJob, generation_uuid)
                if generation_job:
                    mark_generation_failed(
                        db,
                        generation_job,
                        "soft_time_limit_exceeded",
                        "Video generation exceeded the worker time limit. "
                        "Please retry or reduce the number of slides.",
                    )
            logger.exception("Video generation job %s exceeded soft time limit", generation_job_id)
            raise
        except PipelineError as exc:
            with worker_db_session() as db:
                generation_job = db.get(GenerationJob, generation_uuid)
                if generation_job:
                    mark_generation_failed(
                        db,
                        generation_job,
                        exc.code,
                        exc.message,
                        current_step=exc.stage,
                        current_slide=exc.slide_index,
                    )
            logger.exception(
                "Generation job %s: failed at slide %s during stage %s: %s",
                generation_job_id,
                exc.slide_index,
                exc.stage,
                exc.message,
            )
            raise
        except Exception as exc:
            with worker_db_session() as db:
                generation_job = db.get(GenerationJob, generation_uuid)
                if generation_job:
                    mark_generation_failed(
                        db,
                        generation_job,
                        "unexpected_generation_error",
                        str(exc),
                    )
            logger.exception("Video generation job %s failed: %s", generation_job_id, exc)
            raise


generate_video = celery_app.register_task(GenerateVideoTask())


class GenerationAssetError(RuntimeError):
    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code
        self.message = message


def _validate_slide_assets(
    index: int,
    slide_image: bytes,
    audio_bytes: bytes,
    audio_duration: float,
    avatar_clip: bytes,
    avatar_duration: float,
) -> None:
    if not slide_image:
        raise GenerationAssetError(
            "MISSING_SLIDE_RENDER",
            f"Slide render for slide {index} is missing",
        )
    if not audio_bytes or audio_duration <= 0:
        raise GenerationAssetError(
            "MISSING_AUDIO_ASSET",
            f"Audio asset for slide {index} is missing or empty",
        )
    if not avatar_clip or avatar_duration <= 0:
        raise GenerationAssetError(
            "MISSING_AVATAR_CLIP",
            f"Avatar clip for slide {index} is missing or empty",
        )


def _probe_media_duration(media_bytes: bytes, suffix: str) -> float:
    if not media_bytes:
        return 0.0
    with tempfile.TemporaryDirectory() as tmp:
        media_path = Path(tmp) / f"asset{suffix}"
        media_path.write_bytes(media_bytes)
        try:
            result = subprocess.run(
                [
                    "ffprobe",
                    "-v",
                    "error",
                    "-show_entries",
                    "format=duration",
                    "-of",
                    "default=noprint_wrappers=1:nokey=1",
                    str(media_path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=30,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
            return 0.0
    try:
        return float(result.stdout.strip())
    except ValueError:
        return 0.0


def _probe_media_streams(media_bytes: bytes, suffix: str) -> dict[str, bool]:
    streams = {"video": False, "audio": False}
    if not media_bytes:
        return streams
    with tempfile.TemporaryDirectory() as tmp:
        media_path = Path(tmp) / f"asset{suffix}"
        media_path.write_bytes(media_bytes)
        try:
            result = subprocess.run(
                [
                    "ffprobe",
                    "-v",
                    "error",
                    "-show_entries",
                    "stream=codec_type",
                    "-of",
                    "csv=p=0",
                    str(media_path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=30,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
            return streams
    for line in result.stdout.splitlines():
        codec_type = line.strip()
        if codec_type in streams:
            streams[codec_type] = True
    return streams


def _probe_media_info(media_bytes: bytes, suffix: str) -> dict:
    info = {
        "duration_seconds": 0.0,
        "has_video": False,
        "has_audio": False,
        "video_codec": None,
        "audio_codec": None,
        "width": None,
        "height": None,
    }
    if not media_bytes:
        return info
    with tempfile.TemporaryDirectory() as tmp:
        media_path = Path(tmp) / f"asset{suffix}"
        media_path.write_bytes(media_bytes)
        try:
            result = subprocess.run(
                [
                    "ffprobe",
                    "-v",
                    "error",
                    "-show_entries",
                    "format=duration:stream=codec_type,codec_name,width,height",
                    "-of",
                    "json",
                    str(media_path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=30,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
            return info
    try:
        payload = json.loads(result.stdout)
    except json.JSONDecodeError:
        return info
    try:
        info["duration_seconds"] = float((payload.get("format") or {}).get("duration") or 0)
    except (TypeError, ValueError):
        info["duration_seconds"] = 0.0
    for stream in payload.get("streams") or []:
        if not isinstance(stream, dict):
            continue
        if stream.get("codec_type") == "video" and not info["has_video"]:
            info["has_video"] = True
            info["video_codec"] = stream.get("codec_name")
            info["width"] = stream.get("width")
            info["height"] = stream.get("height")
        if stream.get("codec_type") == "audio" and not info["has_audio"]:
            info["has_audio"] = True
            info["audio_codec"] = stream.get("codec_name")
    return info


def _generate_avatar_with_motion_fallback(
    avatar_provider,
    audio_bytes: bytes,
    duration: float,
    avatar_source_url: str,
    wavespeed_api_key: str,
    audio_url: str,
    slide_index: int,
    generation_uuid: uuid.UUID,
) -> tuple[bytes, dict]:
    primary_model = settings.DEFAULT_LIPSYNC_MODEL
    clip = avatar_provider.generate_avatar_clip(
        audio_bytes=audio_bytes,
        duration_seconds=duration,
        avatar_id=avatar_source_url,
        api_key=wavespeed_api_key,
        audio_url=audio_url,
        avatar_source_url=avatar_source_url,
        model_name=primary_model,
        prompt=settings.LIPSYNC_PROMPT,
    )
    motion = _analyze_video_motion(clip)
    metadata = {
        "model": primary_model,
        "fallback_used": False,
        "motion_analysis": motion,
        "prompt": settings.LIPSYNC_PROMPT,
    }
    fallback_model = (settings.FALLBACK_LIPSYNC_MODEL or "").strip()
    if motion.get("almost_static") and fallback_model:
        logger.warning(
            "Primary WaveSpeed lipsync result appears almost static; retrying fallback model: "
            "generation_job=%s slide=%s primary_model=%s fallback_model=%s score=%s",
            generation_uuid,
            slide_index,
            primary_model,
            fallback_model,
            motion.get("motion_score"),
        )
        fallback_clip = avatar_provider.generate_avatar_clip(
            audio_bytes=audio_bytes,
            duration_seconds=duration,
            avatar_id=avatar_source_url,
            api_key=wavespeed_api_key,
            audio_url=audio_url,
            avatar_source_url=avatar_source_url,
            model_name=fallback_model,
            prompt=settings.LIPSYNC_PROMPT,
        )
        fallback_motion = _analyze_video_motion(fallback_clip)
        if float(fallback_motion.get("motion_score") or 0) > float(
            motion.get("motion_score") or 0
        ):
            return fallback_clip, {
                "model": fallback_model,
                "primary_model": primary_model,
                "fallback_used": True,
                "primary_motion_analysis": motion,
                "motion_analysis": fallback_motion,
                "prompt": settings.LIPSYNC_PROMPT,
            }
        metadata["fallback_attempted"] = True
        metadata["fallback_model"] = fallback_model
        metadata["fallback_motion_analysis"] = fallback_motion
    return clip, metadata


def _analyze_video_motion(media_bytes: bytes) -> dict:
    info = _probe_media_info(media_bytes, ".mp4")
    duration = float(info.get("duration_seconds") or 0)
    if not media_bytes or duration <= 0 or not info.get("has_video"):
        return {"motion_score": 0.0, "almost_static": True, "sample_count": 0}
    sample_times = [0.15, max(0.15, duration / 2), max(0.15, duration - 0.15)]
    frames: list[bytes] = []
    with tempfile.TemporaryDirectory() as tmp:
        media_path = Path(tmp) / "avatar.mp4"
        media_path.write_bytes(media_bytes)
        for sample_time in sample_times:
            try:
                result = subprocess.run(
                    [
                        "ffmpeg",
                        "-v",
                        "error",
                        "-ss",
                        f"{sample_time:.3f}",
                        "-i",
                        str(media_path),
                        "-frames:v",
                        "1",
                        "-vf",
                        "scale=64:64,format=gray",
                        "-f",
                        "rawvideo",
                        "-",
                    ],
                    check=True,
                    capture_output=True,
                    timeout=30,
                )
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
                continue
            if result.stdout:
                frames.append(result.stdout)
    if len(frames) < 2:
        return {"motion_score": 0.0, "almost_static": True, "sample_count": len(frames)}
    scores = [_mean_abs_frame_diff(frames[index], frames[index + 1]) for index in range(len(frames) - 1)]
    motion_score = max(scores) if scores else 0.0
    return {
        "motion_score": round(motion_score, 4),
        "almost_static": motion_score < 1.5,
        "sample_count": len(frames),
        "threshold": 1.5,
    }


def _mean_abs_frame_diff(left: bytes, right: bytes) -> float:
    length = min(len(left), len(right))
    if length <= 0:
        return 0.0
    return sum(abs(left[index] - right[index]) for index in range(length)) / length


def _ffmpeg_error_message(exc: subprocess.CalledProcessError) -> str:
    stderr = exc.stderr
    if isinstance(stderr, bytes):
        return stderr.decode("utf-8", errors="ignore")[-1000:] or "FFmpeg failed"
    return str(stderr or "FFmpeg failed")[-1000:]


def _avatar_source_url(settings_row: VideoGenerationSettings, storage) -> str | None:
    if settings_row.avatar_source_url:
        _validate_external_provider_url(settings_row.avatar_source_url)
        return settings_row.avatar_source_url
    if settings_row.avatar_source_asset_id:
        with worker_db_session() as db:
            asset = db.get(Asset, settings_row.avatar_source_asset_id)
            if asset:
                try:
                    if not storage.download_bytes(asset.storage_key):
                        raise GenerationAssetError(
                            "MISSING_AVATAR_ASSET",
                            "Please upload an avatar image before generating the video.",
                        )
                    return _external_asset_url(storage, asset.storage_key)
                except Exception as exc:
                    if isinstance(exc, GenerationAssetError):
                        raise
                    raise GenerationAssetError(
                        "AVATAR_SIGNED_URL_FAILED",
                        "Could not create a signed URL for the avatar asset.",
                    ) from exc
    if settings.DEBUG_AVATAR_SOURCE_URL and not settings.is_production:
        _validate_external_provider_url(settings.DEBUG_AVATAR_SOURCE_URL)
        return settings.DEBUG_AVATAR_SOURCE_URL
    return None


def _external_asset_url(storage, storage_key: str) -> str:
    url = storage.generate_external_download_url(storage_key).url
    _validate_external_provider_url(url)
    return url


def _validate_external_provider_url(url: str) -> None:
    parsed = urlparse(url)
    hostname = parsed.hostname or ""
    if parsed.scheme not in {"http", "https"} or not hostname:
        raise GenerationAssetError(
            "EXTERNAL_ASSET_URL_NOT_PUBLIC",
            "WaveSpeed requires a public URL for audio/avatar assets. Configure a public tunnel or external storage.",
        )
    lowered_host = hostname.lower()
    blocked_hosts = {"localhost", "127.0.0.1", "0.0.0.0", "minio"}
    is_blocked = lowered_host in blocked_hosts or lowered_host.endswith(".local")
    try:
        ip_address = ipaddress.ip_address(lowered_host)
        is_blocked = (
            is_blocked
            or ip_address.is_private
            or ip_address.is_loopback
            or ip_address.is_link_local
        )
    except ValueError:
        pass
    if is_blocked:
        raise GenerationAssetError(
            "EXTERNAL_ASSET_URL_NOT_PUBLIC",
            "WaveSpeed requires a public URL for audio/avatar assets. Configure a public tunnel or external storage.",
        )


def _log_storage_config_for_generation() -> None:
    internal = settings.MINIO_INTERNAL_ENDPOINT or settings.MINIO_ENDPOINT
    public = settings.MINIO_PUBLIC_ENDPOINT or settings.MINIO_ENDPOINT
    logger.info(
        "Storage config for generation: backend=%s internal_host=%s public_host=%s "
        "azure_public_host=%s external_provider_asset_base_url_configured=%s",
        settings.STORAGE_BACKEND,
        _endpoint_host(internal),
        _endpoint_host(public),
        _endpoint_host(settings.AZURE_STORAGE_PUBLIC_BASE_URL),
        bool(settings.EXTERNAL_PROVIDER_ASSET_BASE_URL),
    )


def _endpoint_host(endpoint: str | None) -> str | None:
    if not endpoint:
        return None
    parsed = urlparse(endpoint)
    if parsed.scheme in {"http", "https"}:
        return parsed.netloc
    return endpoint.split("/", 1)[0]


def _avatar_overlay_from_metadata(metadata: dict, resolution: str) -> dict[str, int]:
    output_width, output_height = _resolution_size_for_generation(resolution)
    canvas = metadata.get("canvas") if isinstance(metadata, dict) else None
    canvas_width = 960.0
    canvas_height = 540.0
    avatar: dict | None = None
    if isinstance(canvas, dict):
        canvas_width = float(canvas.get("width") or canvas_width)
        canvas_height = float(canvas.get("height") or canvas_height)
        if isinstance(canvas.get("avatar"), dict):
            avatar = canvas["avatar"]
        else:
            elements = canvas.get("elements") or []
            avatar = next(
                (
                    element
                    for element in elements
                    if isinstance(element, dict) and element.get("type") == "avatar"
                ),
                None,
            )

    default_width = 400
    default_height = 400
    if not avatar:
        return {
            "x": output_width - default_width - 80,
            "y": output_height - default_height - 60,
            "width": default_width,
            "height": default_height,
        }

    scale_x = output_width / max(canvas_width, 1.0)
    scale_y = output_height / max(canvas_height, 1.0)
    width = max(1, int(float(avatar.get("width") or default_width) * scale_x))
    height = max(1, int(float(avatar.get("height") or default_height) * scale_y))
    x = int(float(avatar.get("x") or 0) * scale_x)
    y = int(float(avatar.get("y") or 0) * scale_y)
    return {
        "x": max(0, min(x, output_width - width)),
        "y": max(0, min(y, output_height - height)),
        "width": min(width, output_width),
        "height": min(height, output_height),
    }


def _resolution_size_for_generation(resolution: str) -> tuple[int, int]:
    if resolution == "720p":
        return 1280, 720
    if resolution == "1080p":
        return 1920, 1080
    if "x" in resolution:
        left, right = resolution.lower().split("x", 1)
        return int(left), int(right)
    return 1920, 1080


# ── Función helper para encolar jobs ─────────────────────────────────────────

def enqueue_parse_presentation(job_id: uuid.UUID, presentation_id: uuid.UUID) -> str:
    """Encola un job de parse_presentation y retorna el celery_task_id."""
    result = parse_presentation.apply_async(
        kwargs={
            "job_id": str(job_id),
            "presentation_id": str(presentation_id),
        },
        queue="presentations",
    )
    return result.id


def enqueue_generate_video(
    job_id: uuid.UUID,
    generation_job_id: uuid.UUID,
    project_id: uuid.UUID,
) -> str:
    result = generate_video.apply_async(
        kwargs={
            "job_id": str(job_id),
            "generation_job_id": str(generation_job_id),
            "project_id": str(project_id),
        },
        queue="generation",
    )
    return result.id


def _update_generation_job(
    generation_job_id: uuid.UUID,
    status: str,
    progress: float,
    current_step: str,
    current_slide: int | None = None,
    total_slides: int | None = None,
) -> None:
    with worker_db_session() as db:
        generation_job = db.get(GenerationJob, generation_job_id)
        if generation_job:
            generation_job.status = status
            generation_job.progress_percentage = progress
            generation_job.current_step = current_step
            if current_slide is not None:
                generation_job.current_slide = current_slide
            if total_slides is not None:
                generation_job.total_slides = total_slides
            db.commit()


def _complete_generation_job(
    generation_job_id: uuid.UUID,
    final_asset_id: uuid.UUID,
    result: dict,
) -> None:
    with worker_db_session() as db:
        generation_job = db.get(GenerationJob, generation_job_id)
        if generation_job:
            generation_job.status = "completed"
            generation_job.progress_percentage = 100.0
            generation_job.current_step = "Completed"
            generation_job.final_asset_id = final_asset_id
            generation_job.result = result
            generation_job.completed_at = datetime.now(UTC)
            db.commit()


def _fail_generation_job(
    generation_job_id: uuid.UUID,
    error_code: str,
    error_message: str,
) -> None:
    with worker_db_session() as db:
        generation_job = db.get(GenerationJob, generation_job_id)
        if generation_job:
            generation_job.status = "failed"
            generation_job.error_code = error_code
            generation_job.error_message = error_message[:2000]
            generation_job.completed_at = datetime.now(UTC)
            db.commit()


def _create_asset(
    organization_id: uuid.UUID,
    project_id: uuid.UUID,
    slide_id: uuid.UUID | None,
    asset_type: str,
    storage_key: str,
    filename: str,
    mime_type: str,
    size_bytes: int,
    duration_seconds: float | None = None,
    metadata_json: dict | None = None,
) -> Asset:
    with worker_db_session() as db:
        asset = Asset(
            organization_id=organization_id,
            project_id=project_id,
            slide_id=slide_id,
            asset_type=asset_type,
            storage_key=storage_key,
            filename=filename,
            mime_type=mime_type,
            size_bytes=size_bytes,
            duration_seconds=duration_seconds,
            metadata_json=metadata_json,
        )
        db.add(asset)
        db.commit()
        db.refresh(asset)
        return asset
