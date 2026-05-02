import io
import logging
import shutil
import subprocess
import tempfile
import uuid
from pathlib import Path

from pptx import Presentation as PptxPresentation

logger = logging.getLogger(__name__)


def render_slide_previews(
    pptx_bytes: bytes,
    presentation_id: uuid.UUID,
    original_filename: str,
    storage,
) -> dict[int, str]:
    """Render PPT/PPTX slides to PNG previews and store them in object storage.

    Rendering is best-effort: text parsing remains useful in test/CI environments where
    LibreOffice or Poppler may not be installed.
    """
    return _render_slide_images(
        pptx_bytes=pptx_bytes,
        presentation_id=presentation_id,
        original_filename=original_filename,
        storage=storage,
        folder="previews",
        label="preview",
    )


def render_slide_backgrounds(
    pptx_bytes: bytes,
    presentation_id: uuid.UUID,
    original_filename: str,
    storage,
) -> dict[int, str]:
    """Render slide PNG backgrounds with editable text removed.

    The editor uses these images as non-editable layers for complex PPT content
    while rendering detected text as structured editable elements.
    """
    return _render_slide_images(
        pptx_bytes=_remove_text_from_pptx(pptx_bytes),
        presentation_id=presentation_id,
        original_filename=original_filename,
        storage=storage,
        folder="backgrounds",
        label="background",
    )


def _render_slide_images(
    pptx_bytes: bytes,
    presentation_id: uuid.UUID,
    original_filename: str,
    storage,
    folder: str,
    label: str,
) -> dict[int, str]:
    if not shutil.which("soffice") or not shutil.which("pdftoppm"):
        logger.warning(
            "Slide %s rendering skipped: LibreOffice soffice or pdftoppm is unavailable.",
            label,
        )
        return {}

    suffix = Path(original_filename).suffix.lower()
    if suffix not in {".ppt", ".pptx"}:
        suffix = ".pptx"

    try:
        with tempfile.TemporaryDirectory() as tmp:
            tmpdir = Path(tmp)
            input_path = tmpdir / f"presentation{suffix}"
            input_path.write_bytes(pptx_bytes)

            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(tmpdir),
                    str(input_path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=120,
            )

            pdf_path = tmpdir / "presentation.pdf"
            if not pdf_path.exists():
                pdfs = list(tmpdir.glob("*.pdf"))
                if not pdfs:
                    raise RuntimeError("LibreOffice did not produce a PDF preview source")
                pdf_path = pdfs[0]

            output_prefix = tmpdir / "slide"
            subprocess.run(
                [
                    "pdftoppm",
                    "-png",
                    "-r",
                    "144",
                    str(pdf_path),
                    str(output_prefix),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=120,
            )

            preview_keys: dict[int, str] = {}
            image_paths = sorted(tmpdir.glob("slide-*.png"), key=_rendered_slide_number)
            version = uuid.uuid4().hex
            for index, image_path in enumerate(image_paths, 1):
                key = f"presentations/{presentation_id}/{folder}/{version}/slide-{index}.png"
                storage.upload_file(
                    key=key,
                    data=image_path.read_bytes(),
                    content_type="image/png",
                )
                preview_keys[index] = key
                logger.info(
                    "Rendered slide %s %s for presentation %s to %s",
                    index,
                    label,
                    presentation_id,
                    key,
                )

            logger.info(
                "Rendered %s slide %ss for presentation %s",
                len(preview_keys),
                label,
                presentation_id,
            )
            return preview_keys
    except Exception as exc:
        logger.warning("Slide %s rendering failed: %s", label, exc)
        return {}


def _remove_text_from_pptx(pptx_bytes: bytes) -> bytes:
    deck = PptxPresentation(io.BytesIO(pptx_bytes))
    for slide in deck.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text = ""

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _rendered_slide_number(path: Path) -> int:
    try:
        return int(path.stem.rsplit("-", 1)[-1])
    except ValueError:
        return 0
