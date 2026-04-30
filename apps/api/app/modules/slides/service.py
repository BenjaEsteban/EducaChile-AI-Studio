import io
import uuid

from fastapi import HTTPException, status
from pptx import Presentation as PptxPresentation

from app.modules.presentations.rendering import render_slide_previews
from app.modules.projects.models import Slide
from app.modules.projects.service import MOCK_ORG_ID
from app.modules.slides.repository import SlideRepository
from app.modules.slides.schemas import SlideRead, SlideUpdate
from app.providers.storage import get_storage


class SlideService:
    def __init__(self, repo: SlideRepository) -> None:
        self.repo = repo

    def list_by_presentation(self, presentation_id: uuid.UUID) -> list[SlideRead]:
        slides = self.repo.list_by_presentation(presentation_id, MOCK_ORG_ID)
        if slides is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Presentation not found",
            )
        return [SlideRead.from_model(slide) for slide in slides]

    def get_or_404(self, slide_id: uuid.UUID) -> SlideRead:
        slide = self.repo.get_by_id(slide_id, MOCK_ORG_ID)
        if not slide:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Slide not found",
            )
        return SlideRead.from_model(slide)

    def update(self, slide_id: uuid.UUID, data: SlideUpdate) -> SlideRead:
        slide = self.repo.get_by_id(slide_id, MOCK_ORG_ID)
        if not slide:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Slide not found",
            )

        update_data = data.model_dump(exclude_unset=True)
        if "title" in update_data:
            slide.title = update_data["title"]
        if "notes" in update_data:
            slide.notes = update_data["notes"]

        metadata = dict(slide.metadata_ or {})
        if "metadata" in update_data and update_data["metadata"] is not None:
            metadata.update(update_data["metadata"])
        if "dialogue" in update_data:
            metadata["dialogue"] = update_data["dialogue"] or ""
        if "visible_text" in update_data:
            metadata["visible_text"] = update_data["visible_text"] or ""
        slide.metadata_ = metadata

        has_text_update = "title" in update_data or "visible_text" in update_data
        if has_text_update and _contains_canvas_text_blocks(metadata):
            self._update_pptx_text_and_previews(slide)

        return SlideRead.from_model(self.repo.save(slide))

    def _update_pptx_text_and_previews(self, slide: Slide) -> None:
        presentation = slide.presentation
        if not presentation:
            return

        storage = get_storage()
        pptx_bytes = storage.download_file(presentation.storage_key)
        deck = PptxPresentation(io.BytesIO(pptx_bytes))
        slide_index = slide.position - 1
        if slide_index < 0 or slide_index >= len(deck.slides):
            return

        text_blocks = _get_canvas_text_blocks(slide.metadata_ or {})
        ppt_slide = deck.slides[slide_index]
        for block in text_blocks:
            shape_index = _shape_index_from_block_id(str(block.get("id") or ""))
            text = str(block.get("text") or "")
            if shape_index is None or shape_index >= len(ppt_slide.shapes):
                continue
            shape = ppt_slide.shapes[shape_index]
            if getattr(shape, "has_text_frame", False):
                shape.text = text

        buffer = io.BytesIO()
        deck.save(buffer)
        updated_pptx = buffer.getvalue()
        storage.upload_file(
            key=presentation.storage_key,
            data=updated_pptx,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        preview_keys = render_slide_previews(
            pptx_bytes=updated_pptx,
            presentation_id=presentation.id,
            original_filename=presentation.original_filename,
            storage=storage,
        )
        if not preview_keys:
            return

        slides = self.repo.list_by_presentation_id(presentation.id)
        for current_slide in slides:
            preview_key = preview_keys.get(current_slide.position)
            if not preview_key:
                continue
            current_slide.thumbnail_key = preview_key
            current_metadata = dict(current_slide.metadata_ or {})
            current_metadata["rendered_image_key"] = preview_key
            current_slide.metadata_ = current_metadata
        self.repo.commit()


def _contains_canvas_text_blocks(metadata: dict) -> bool:
    return bool(_get_canvas_text_blocks(metadata))


def _get_canvas_text_blocks(metadata: dict) -> list[dict]:
    canvas = metadata.get("canvas")
    if not isinstance(canvas, dict):
        return []
    text_blocks = canvas.get("text_blocks")
    if not isinstance(text_blocks, list):
        return []
    return [block for block in text_blocks if isinstance(block, dict)]


def _shape_index_from_block_id(block_id: str) -> int | None:
    try:
        return int(block_id.rsplit("-", 1)[-1])
    except ValueError:
        return None
